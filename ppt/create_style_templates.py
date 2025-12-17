"""创建多样化高质量 PPT 模板

每个模板使用完全不同的布局方式：
1. 斜切分割 - 对角线切割
2. Bento网格 - 日式便当盒布局
3. 卡片堆叠 - 重叠层次效果
4. 大字报风 - 超大文字主导
5. 杂志排版 - 多栏混排
6. 几何拼接 - 三角形/多边形
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml


TEMPLATES_DIR = Path(__file__).parent / "pptx_templates"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_shape(slide, shape_type, left, top, width, height, fill_color=None,
              line_color=None, line_width=0, rotation=0):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color and line_width > 0:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    if rotation:
        shape.rotation = rotation
    return shape


def add_text_box(slide, left, top, width, height, text="", font_size=14,
                 font_color=None, bold=False, align=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    """添加文本框"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    if font_color:
        p.font.color.rgb = font_color
    return box


def move_to_back(slide, shape):
    """移到底层"""
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


# ==================== 模板 1: 斜切分割风格 ====================

def create_diagonal_split_template():
    """斜切分割模板 - 对角线切割页面，动感现代"""
    prs = create_presentation()

    # 配色：深紫 + 珊瑚粉 + 白
    c_dark = RGBColor(45, 27, 62)
    c_accent = RGBColor(255, 107, 107)
    c_secondary = RGBColor(78, 205, 196)
    c_light = RGBColor(250, 250, 250)
    c_white = RGBColor(255, 255, 255)

    # ===== 封面：斜切两半 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    # 左下深色三角（使用大矩形旋转模拟）
    # 创建一个超大的旋转矩形来模拟斜切效果
    diag1 = add_shape(slide1, MSO_SHAPE.RIGHT_TRIANGLE,
                      Inches(-2), Inches(-1), Inches(16), Inches(10), c_dark)

    # 右上浅色背景
    bg = add_shape(slide1, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_light)
    move_to_back(slide1, bg)

    # 斜线装饰
    line1 = add_shape(slide1, MSO_SHAPE.RECTANGLE,
                      Inches(6), Inches(0), Pt(4), Inches(8), c_accent, rotation=-25)
    line2 = add_shape(slide1, MSO_SHAPE.RECTANGLE,
                      Inches(6.5), Inches(-0.5), Pt(2), Inches(8), c_secondary, rotation=-25)

    # 标题（在深色区域）
    add_text_box(slide1, Inches(0.8), Inches(4), Inches(5), Inches(1.2),
                 "斜切", 72, c_white, bold=True)
    add_text_box(slide1, Inches(0.8), Inches(5.3), Inches(5), Inches(0.8),
                 "DIAGONAL", 36, c_accent)

    # 右侧副标题
    add_text_box(slide1, Inches(8), Inches(1.5), Inches(4.5), Inches(1),
                 "创意设计模板", 28, c_dark, bold=True)
    add_text_box(slide1, Inches(8), Inches(2.5), Inches(4.5), Inches(0.6),
                 "打破常规 · 动感现代", 16, RGBColor(100, 100, 100))

    # ===== 内容页：斜切布局 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    bg2 = add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_light)
    move_to_back(slide2, bg2)

    # 左侧斜切色块
    add_shape(slide2, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(-1), Inches(-1), Inches(6), Inches(9), c_dark)

    # 顶部装饰条（斜向）
    add_shape(slide2, MSO_SHAPE.PARALLELOGRAM,
              Inches(4), Inches(-0.5), Inches(10), Inches(1.2), c_accent)

    # 标题区
    add_text_box(slide2, Inches(0.6), Inches(0.8), Inches(3.5), Inches(0.9),
                 "01", 48, c_white, bold=True)
    add_text_box(slide2, Inches(0.6), Inches(1.8), Inches(3.5), Inches(0.7),
                 "章节标题", 24, c_white, bold=True)

    # 右侧内容区
    add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(5), Inches(1.5), Inches(7.5), Inches(5.3), c_white)
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              Inches(5), Inches(1.5), Pt(5), Inches(5.3), c_secondary)

    # ===== 结束页 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape(slide3, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(7), Inches(-1), Inches(8), Inches(9), c_accent)
    bg3 = add_shape(slide3, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_dark)
    move_to_back(slide3, bg3)

    add_text_box(slide3, Inches(1), Inches(2.5), Inches(6), Inches(1.5),
                 "THANKS", 80, c_white, bold=True)
    add_text_box(slide3, Inches(1), Inches(4.2), Inches(6), Inches(0.6),
                 "感谢观看", 24, c_secondary)

    prs.save(str(TEMPLATES_DIR / "style_diagonal_split.pptx"))
    print("✓ 斜切分割模板: style_diagonal_split.pptx")


# ==================== 模板 2: Bento网格风格 ====================

def create_bento_grid_template():
    """Bento网格模板 - 日式便当盒布局，信息密集但有序"""
    prs = create_presentation()

    # 配色：深灰 + 亮黄 + 青色
    c_dark = RGBColor(30, 30, 35)
    c_yellow = RGBColor(255, 214, 0)
    c_cyan = RGBColor(0, 210, 211)
    c_magenta = RGBColor(255, 46, 99)
    c_light = RGBColor(245, 245, 245)
    c_white = RGBColor(255, 255, 255)

    # ===== 封面：大网格 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    bg = add_shape(slide1, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_dark)
    move_to_back(slide1, bg)

    gap = Inches(0.15)

    # 左上大格子（标题）
    add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE,
              gap, gap, Inches(7) - gap, Inches(5) - gap, c_yellow)
    add_text_box(slide1, Inches(0.5), Inches(1.5), Inches(6), Inches(2),
                 "BENTO\nGRID", 64, c_dark, bold=True)
    add_text_box(slide1, Inches(0.5), Inches(3.8), Inches(6), Inches(0.6),
                 "网格布局设计", 20, c_dark)

    # 右上格子
    add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(7) + gap, gap, Inches(3) - gap*2, Inches(2.5) - gap, c_cyan)
    add_text_box(slide1, Inches(7.3), Inches(0.6), Inches(2.4), Inches(1.5),
                 "01\n模块化", 24, c_dark, bold=True, align=PP_ALIGN.CENTER)

    # 右中格子
    add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(7) + gap, Inches(2.5) + gap, Inches(3) - gap*2, Inches(2.5) - gap, c_magenta)
    add_text_box(slide1, Inches(7.3), Inches(3.1), Inches(2.4), Inches(1.5),
                 "02\n层次感", 24, c_white, bold=True, align=PP_ALIGN.CENTER)

    # 右侧窄格子
    add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(10) + gap, gap, Inches(3.333) - gap*2, Inches(5) - gap, c_light)
    add_text_box(slide1, Inches(10.3), Inches(2), Inches(2.7), Inches(1),
                 "2024", 28, c_dark, bold=True, align=PP_ALIGN.CENTER)

    # 底部长条
    add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE,
              gap, Inches(5) + gap, Inches(13.333) - gap*2, Inches(2.5) - gap*2, c_white)
    add_text_box(slide1, Inches(0.5), Inches(5.8), Inches(12), Inches(0.8),
                 "信息密集但井井有条 · 适合数据展示与功能介绍", 18, c_dark, align=PP_ALIGN.CENTER)

    # ===== 内容页：四格布局 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    bg2 = add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_light)
    move_to_back(slide2, bg2)

    # 顶部标题条
    add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1), c_dark)
    add_text_box(slide2, Inches(0.5), Inches(0.25), Inches(5), Inches(0.6),
                 "内容概览", 28, c_white, bold=True)

    # 四个内容格子
    colors = [c_yellow, c_cyan, c_magenta, c_white]
    text_colors = [c_dark, c_dark, c_white, c_dark]
    labels = ["功能一", "功能二", "功能三", "功能四"]

    for i in range(4):
        row, col = i // 2, i % 2
        x = Inches(0.3) + col * Inches(6.5)
        y = Inches(1.3) + row * Inches(3)
        add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                  x, y, Inches(6.2), Inches(2.8), colors[i])
        add_text_box(slide2, x + Inches(0.3), y + Inches(0.3), Inches(5.6), Inches(0.8),
                     f"0{i+1} {labels[i]}", 22, text_colors[i], bold=True)

    # ===== 结束页 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    bg3 = add_shape(slide3, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_dark)
    move_to_back(slide3, bg3)

    # 中央大格子
    add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(3), Inches(1.5), Inches(7.333), Inches(4.5), c_yellow)
    add_text_box(slide3, Inches(3.5), Inches(2.8), Inches(6.333), Inches(1.5),
                 "THANK\nYOU", 56, c_dark, bold=True, align=PP_ALIGN.CENTER)

    # 小装饰格子
    add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(0.3), Inches(0.3), Inches(2), Inches(2), c_cyan)
    add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(11), Inches(5.2), Inches(2), Inches(2), c_magenta)

    prs.save(str(TEMPLATES_DIR / "style_bento_grid.pptx"))
    print("✓ Bento网格模板: style_bento_grid.pptx")


# ==================== 模板 3: 卡片堆叠风格 ====================

def create_card_stack_template():
    """卡片堆叠模板 - 重叠层次效果，有深度感"""
    prs = create_presentation()

    # 配色：深蓝渐变 + 白卡片
    c_bg_dark = RGBColor(15, 23, 42)
    c_bg_mid = RGBColor(30, 41, 59)
    c_accent = RGBColor(99, 102, 241)  # 靛蓝
    c_accent2 = RGBColor(236, 72, 153)  # 粉
    c_white = RGBColor(255, 255, 255)
    c_card = RGBColor(248, 250, 252)
    c_text = RGBColor(51, 65, 85)

    # ===== 封面：堆叠卡片 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    bg = add_shape(slide1, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_bg_dark)
    move_to_back(slide1, bg)

    # 背景装饰圆（模拟光晕）
    add_shape(slide1, MSO_SHAPE.OVAL,
              Inches(6), Inches(-2), Inches(10), Inches(10), c_bg_mid)

    # 堆叠的卡片（从后到前）
    # 卡片3（最后面）
    add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(3), Inches(1.2), Inches(7), Inches(5), c_accent, rotation=8)

    # 卡片2
    add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(3.3), Inches(1), Inches(7), Inches(5), c_accent2, rotation=4)

    # 卡片1（最前面，主卡片）
    main_card = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(3.5), Inches(0.8), Inches(7), Inches(5), c_white)

    add_text_box(slide1, Inches(4.2), Inches(2), Inches(5.5), Inches(1.5),
                 "CARD\nSTACK", 52, c_bg_dark, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(4.2), Inches(4), Inches(5.5), Inches(0.6),
                 "层叠卡片设计", 20, c_text, align=PP_ALIGN.CENTER)

    # 左下角文字
    add_text_box(slide1, Inches(0.5), Inches(6), Inches(3), Inches(0.8),
                 "有层次 · 有深度 · 有重点", 14, RGBColor(148, 163, 184))

    # ===== 内容页：卡片布局 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    bg2 = add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_card)
    move_to_back(slide2, bg2)

    # 顶部色条
    add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.08), c_accent)

    # 标题
    add_text_box(slide2, Inches(0.6), Inches(0.4), Inches(8), Inches(0.8),
                 "内容标题", 32, c_bg_dark, bold=True)

    # 三张堆叠小卡片
    for i in range(3):
        offset = i * 0.15
        shade = 230 - i * 15
        c = RGBColor(shade, shade, shade) if i > 0 else c_white
        add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                  Inches(0.6 + offset), Inches(1.5 + offset),
                  Inches(5.5), Inches(5), c)

    # 右侧内容卡片
    add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(6.8), Inches(1.5), Inches(6), Inches(5), c_white)
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              Inches(6.8), Inches(1.5), Inches(6), Inches(0.8), c_accent)
    add_text_box(slide2, Inches(7), Inches(1.65), Inches(5.5), Inches(0.5),
                 "要点内容", 18, c_white, bold=True)

    # ===== 结束页 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    bg3 = add_shape(slide3, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_bg_dark)
    move_to_back(slide3, bg3)

    # 堆叠效果
    add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(4.2), Inches(1.8), Inches(5), Inches(4), c_accent, rotation=6)
    add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(4.4), Inches(1.6), Inches(5), Inches(4), c_accent2, rotation=3)
    add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(4.6), Inches(1.4), Inches(5), Inches(4), c_white)

    add_text_box(slide3, Inches(5), Inches(2.5), Inches(4), Inches(1.5),
                 "Thanks", 48, c_bg_dark, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide3, Inches(0), Inches(6.5), SLIDE_WIDTH, Inches(0.5),
                 "期待与您的下次交流", 16, RGBColor(148, 163, 184), align=PP_ALIGN.CENTER)

    prs.save(str(TEMPLATES_DIR / "style_card_stack.pptx"))
    print("✓ 卡片堆叠模板: style_card_stack.pptx")


# ==================== 模板 4: 大字报风格 ====================

def create_bold_typography_template():
    """大字报风格 - 超大文字主导，冲击力强"""
    prs = create_presentation()

    # 配色：黑白为主 + 亮色点缀
    c_black = RGBColor(0, 0, 0)
    c_white = RGBColor(255, 255, 255)
    c_accent = RGBColor(255, 59, 48)  # 苹果红
    c_gray = RGBColor(142, 142, 147)

    # ===== 封面：超大文字 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    bg = add_shape(slide1, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_white)
    move_to_back(slide1, bg)

    # 超大标题（填满屏幕）
    add_text_box(slide1, Inches(-0.3), Inches(0.5), Inches(14), Inches(3),
                 "BOLD", 180, c_black, bold=True)
    add_text_box(slide1, Inches(-0.3), Inches(3.5), Inches(14), Inches(2.5),
                 "TYPE", 160, c_black, bold=True)

    # 红色装饰块
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              Inches(10), Inches(5.5), Inches(3), Inches(1.5), c_accent)

    # 副标题
    add_text_box(slide1, Inches(0.5), Inches(6.3), Inches(6), Inches(0.6),
                 "大字报风格 · 视觉冲击", 18, c_gray)

    # ===== 内容页：大字+小内容 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    bg2 = add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_black)
    move_to_back(slide2, bg2)

    # 大号章节数字
    add_text_box(slide2, Inches(0.3), Inches(0), Inches(6), Inches(3.5),
                 "01", 200, c_white, bold=True)

    # 章节标题
    add_text_box(slide2, Inches(0.5), Inches(4), Inches(5), Inches(1),
                 "核心观点", 36, c_white, bold=True)

    # 红色下划线
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              Inches(0.5), Inches(5), Inches(2), Pt(6), c_accent)

    # 右侧内容区
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              Inches(7), Inches(1), Inches(6), Inches(5.5), c_white)
    add_text_box(slide2, Inches(7.3), Inches(1.3), Inches(5.4), Inches(5),
                 "在这里写入内容要点\n\n• 要点一\n• 要点二\n• 要点三",
                 18, c_black)

    # ===== 结束页：纯文字 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    bg3 = add_shape(slide3, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_accent)
    move_to_back(slide3, bg3)

    add_text_box(slide3, Inches(0), Inches(2), SLIDE_WIDTH, Inches(2),
                 "THE", 100, c_white, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide3, Inches(0), Inches(3.8), SLIDE_WIDTH, Inches(2),
                 "END", 140, c_white, bold=True, align=PP_ALIGN.CENTER)

    prs.save(str(TEMPLATES_DIR / "style_bold_typography.pptx"))
    print("✓ 大字报风格模板: style_bold_typography.pptx")


# ==================== 模板 5: 杂志排版风格 ====================

def create_magazine_layout_template():
    """杂志排版风格 - 多栏混排，编辑感"""
    prs = create_presentation()

    # 配色：优雅米色系
    c_cream = RGBColor(250, 245, 240)
    c_dark = RGBColor(35, 31, 32)
    c_accent = RGBColor(200, 100, 80)  # 砖红
    c_gold = RGBColor(180, 150, 100)
    c_white = RGBColor(255, 255, 255)

    # ===== 封面：杂志封面 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    bg = add_shape(slide1, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_cream)
    move_to_back(slide1, bg)

    # 左侧大色块（图片占位区）
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              0, 0, Inches(7), SLIDE_HEIGHT, c_dark)

    # 杂志名（竖排感觉）
    add_text_box(slide1, Inches(0.5), Inches(0.5), Inches(2), Inches(0.8),
                 "MAGAZINE", 14, c_gold, bold=True)

    # 大标题（在色块上）
    add_text_box(slide1, Inches(0.5), Inches(2.5), Inches(6), Inches(2),
                 "STYLE", 90, c_white, bold=True)

    # 右侧文字区
    add_text_box(slide1, Inches(7.5), Inches(1), Inches(5), Inches(1.5),
                 "杂志风格", 48, c_dark, bold=True)

    # 装饰线
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              Inches(7.5), Inches(2.6), Inches(3), Pt(2), c_accent)

    # 副标题
    add_text_box(slide1, Inches(7.5), Inches(3), Inches(5), Inches(1.5),
                 "Editorial Design\n编辑设计风格", 18, c_dark)

    # 底部信息
    add_text_box(slide1, Inches(7.5), Inches(6.5), Inches(5), Inches(0.5),
                 "VOL.01 / 2024", 12, c_gold)

    # ===== 内容页：三栏布局 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    bg2 = add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_cream)
    move_to_back(slide2, bg2)

    # 顶部线条
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              Inches(0.5), Inches(0.4), Inches(12.333), Pt(1), c_dark)

    # 大标题
    add_text_box(slide2, Inches(0.5), Inches(0.6), Inches(8), Inches(1),
                 "章节标题", 42, c_dark, bold=True)

    # 三栏布局
    col_width = Inches(3.8)
    for i in range(3):
        x = Inches(0.5) + i * (col_width + Inches(0.3))

        # 栏目标题
        add_shape(slide2, MSO_SHAPE.RECTANGLE,
                  x, Inches(2), col_width, Pt(3), c_accent if i == 0 else c_gold)
        add_text_box(slide2, x, Inches(2.2), col_width, Inches(0.6),
                     f"栏目 0{i+1}", 14, c_dark, bold=True)

        # 栏目内容区
        add_shape(slide2, MSO_SHAPE.RECTANGLE,
                  x, Inches(3), col_width, Inches(4), c_white)

    # ===== 结束页 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # 左右分割
    add_shape(slide3, MSO_SHAPE.RECTANGLE,
              0, 0, SLIDE_WIDTH / 2, SLIDE_HEIGHT, c_dark)
    bg3 = add_shape(slide3, MSO_SHAPE.RECTANGLE,
              SLIDE_WIDTH / 2, 0, SLIDE_WIDTH / 2, SLIDE_HEIGHT, c_cream)
    move_to_back(slide3, bg3)

    add_text_box(slide3, Inches(0.5), Inches(3), Inches(6), Inches(1.5),
                 "MERCI", 72, c_white, bold=True)

    add_text_box(slide3, Inches(7), Inches(3), Inches(5.5), Inches(1),
                 "感谢阅读", 36, c_dark, bold=True)
    add_shape(slide3, MSO_SHAPE.RECTANGLE,
              Inches(7), Inches(4.2), Inches(2), Pt(3), c_accent)

    prs.save(str(TEMPLATES_DIR / "style_magazine_layout.pptx"))
    print("✓ 杂志排版模板: style_magazine_layout.pptx")


# ==================== 模板 6: 几何拼接风格 ====================

def create_geometric_mosaic_template():
    """几何拼接风格 - 三角形/多边形组合，艺术感"""
    prs = create_presentation()

    # 配色：多彩几何
    c_blue = RGBColor(59, 130, 246)
    c_purple = RGBColor(139, 92, 246)
    c_pink = RGBColor(236, 72, 153)
    c_orange = RGBColor(249, 115, 22)
    c_green = RGBColor(34, 197, 94)
    c_dark = RGBColor(30, 30, 35)
    c_light = RGBColor(250, 250, 250)
    c_white = RGBColor(255, 255, 255)

    # ===== 封面：几何拼接 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    bg = add_shape(slide1, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_light)
    move_to_back(slide1, bg)

    # 左侧几何拼接
    add_shape(slide1, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(0), Inches(0), Inches(4), Inches(4), c_blue)
    add_shape(slide1, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(0), Inches(4), Inches(4), Inches(3.5), c_purple, rotation=90)
    add_shape(slide1, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(4), Inches(0), Inches(3), Inches(4), c_pink, rotation=180)

    # 右下几何
    add_shape(slide1, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(10), Inches(4.5), Inches(3.333), Inches(3), c_orange)
    add_shape(slide1, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(11), Inches(5.5), Inches(2.333), Inches(2), c_green, rotation=90)

    # 中央内容
    add_text_box(slide1, Inches(5), Inches(2.5), Inches(7), Inches(1.5),
                 "GEOMETRIC", 56, c_dark, bold=True)
    add_text_box(slide1, Inches(5), Inches(4.2), Inches(7), Inches(0.8),
                 "几何拼接设计", 24, c_dark)

    # 装饰线
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              Inches(5), Inches(5.2), Inches(4), Pt(4), c_blue)

    # ===== 内容页 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    bg2 = add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_white)
    move_to_back(slide2, bg2)

    # 左上角几何装饰
    add_shape(slide2, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(0), Inches(0), Inches(2.5), Inches(2.5), c_blue)
    add_shape(slide2, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(0), Inches(0), Inches(1.5), Inches(1.5), c_purple)

    # 标题
    add_text_box(slide2, Inches(3), Inches(0.5), Inches(8), Inches(0.9),
                 "内容标题", 36, c_dark, bold=True)

    # 彩色指示条
    colors = [c_blue, c_purple, c_pink, c_orange, c_green]
    for i, color in enumerate(colors):
        add_shape(slide2, MSO_SHAPE.RECTANGLE,
                  Inches(3 + i * 0.4), Inches(1.4), Inches(0.3), Pt(5), color)

    # 内容区
    add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(0.5), Inches(2), Inches(12.333), Inches(5), c_light)

    # 右下角装饰
    add_shape(slide2, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(11.333), Inches(5.5), Inches(2), Inches(2), c_pink, rotation=180)

    # ===== 结束页 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    bg3 = add_shape(slide3, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, c_dark)
    move_to_back(slide3, bg3)

    # 多彩几何背景
    add_shape(slide3, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(0), Inches(0), Inches(5), Inches(5), c_blue)
    add_shape(slide3, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(8.333), Inches(2.5), Inches(5), Inches(5), c_purple)
    add_shape(slide3, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(0), Inches(5), Inches(4), Inches(2.5), c_pink, rotation=90)
    add_shape(slide3, MSO_SHAPE.RIGHT_TRIANGLE,
              Inches(10), Inches(0), Inches(3.333), Inches(3), c_orange, rotation=180)

    # 中央文字
    add_text_box(slide3, Inches(0), Inches(2.8), SLIDE_WIDTH, Inches(1.5),
                 "THANKS", 72, c_white, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide3, Inches(0), Inches(4.5), SLIDE_WIDTH, Inches(0.6),
                 "感谢观看 · 多彩世界", 20, c_light, align=PP_ALIGN.CENTER)

    prs.save(str(TEMPLATES_DIR / "style_geometric_mosaic.pptx"))
    print("✓ 几何拼接模板: style_geometric_mosaic.pptx")


# ==================== 主函数 ====================

def create_all_style_templates():
    """创建所有风格化模板"""
    print("\n" + "=" * 60)
    print("创建多样化风格 PPT 模板")
    print("每个模板使用完全不同的布局方式")
    print("=" * 60)

    TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)

    create_diagonal_split_template()    # 斜切分割
    create_bento_grid_template()        # Bento网格
    create_card_stack_template()        # 卡片堆叠
    create_bold_typography_template()   # 大字报
    create_magazine_layout_template()   # 杂志排版
    create_geometric_mosaic_template()  # 几何拼接

    print("=" * 60)
    print("完成！创建了 6 个不同风格的模板：")
    print("  1. 斜切分割 - 对角线动感布局")
    print("  2. Bento网格 - 日式便当盒布局")
    print("  3. 卡片堆叠 - 层次重叠效果")
    print("  4. 大字报风 - 超大文字冲击")
    print("  5. 杂志排版 - 多栏编辑风格")
    print("  6. 几何拼接 - 三角形艺术感")
    print("=" * 60 + "\n")


if __name__ == "__main__":
    create_all_style_templates()
