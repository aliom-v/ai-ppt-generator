"""模板生成脚本 - 创建多种风格的 PPT 模板"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


TEMPLATES_DIR = os.path.join(os.path.dirname(__file__), "pptx_templates")


def create_hackathon_neon_template():
    """创建黑客松霓虹风格模板 - 纯黑背景+霓虹绿+故障紫"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 颜色配置
    bg_black = RGBColor(0, 0, 0)          # 纯黑背景
    grid_gray = RGBColor(30, 30, 30)      # 深灰网格
    neon_green = RGBColor(57, 255, 20)    # 霓虹绿 #39FF14
    glitch_purple = RGBColor(255, 0, 255) # 故障紫 #FF00FF
    text_white = RGBColor(255, 255, 255)
    text_gray = RGBColor(180, 180, 180)
    
    # ========== 封面页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 纯黑背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_black
    bg.line.fill.background()
    
    # 移到最底层
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # 网格线效果 - 横线
    for i in range(1, 15):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(i * 0.5),
            Inches(13.333), Pt(0.5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = grid_gray
        line.line.fill.background()
    
    # 网格线效果 - 竖线
    for i in range(1, 27):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(i * 0.5), Inches(0),
            Pt(0.5), Inches(7.5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = grid_gray
        line.line.fill.background()
    
    # 左侧霓虹绿装饰条
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.15), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = neon_green
    left_bar.line.fill.background()
    
    # 顶部故障紫装饰线
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.08)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = glitch_purple
    top_line.line.fill.background()
    
    # 标题装饰线（霓虹绿）
    title_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(3.0),
        Inches(4), Pt(3)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = neon_green
    title_line.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "// HACKATHON PROJECT"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = neon_green
    p.font.name = "Consolas"
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(0.8))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.text = "> 硬核技术 · 终端美学 · 代码即艺术"
    sp.font.size = Pt(24)
    sp.font.color.rgb = glitch_purple
    sp.font.name = "Consolas"
    
    # 底部装饰 - 模拟终端光标
    cursor_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(5), Inches(0.5))
    ctf = cursor_box.text_frame
    cp = ctf.paragraphs[0]
    cp.text = "$ ./start_presentation.sh █"
    cp.font.size = Pt(16)
    cp.font.color.rgb = neon_green
    cp.font.name = "Consolas"
    
    # ========== 内容页 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 纯黑背景
    bg2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = bg_black
    bg2.line.fill.background()
    
    spTree2 = slide2.shapes._spTree
    sp2 = bg2._element
    spTree2.remove(sp2)
    spTree2.insert(2, sp2)
    
    # 简化网格 - 横线
    for i in range(1, 15):
        line = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(i * 0.5),
            Inches(13.333), Pt(0.3)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = grid_gray
        line.line.fill.background()
    
    # 顶部霓虹绿装饰条
    top_bar = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.1)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = neon_green
    top_bar.line.fill.background()
    
    # 左侧故障紫装饰条
    left_accent = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.08), Inches(7.5)
    )
    left_accent.fill.solid()
    left_accent.fill.fore_color.rgb = glitch_purple
    left_accent.line.fill.background()
    
    # 标题区域背景
    title_bg = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.4),
        Inches(12.333), Inches(1)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(15, 15, 15)
    title_bg.line.fill.background()
    
    # 页面标题
    title_box2 = slide2.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11), Inches(0.8))
    tf2 = title_box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "# 章节标题"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = neon_green
    p2.font.name = "Consolas"
    
    # 内容区域提示
    content_box = slide2.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(11), Inches(4))
    ctf2 = content_box.text_frame
    ctf2.word_wrap = True
    cp2 = ctf2.paragraphs[0]
    cp2.text = "// 在此添加内容\n// 支持代码块、终端截图、网络拓扑图"
    cp2.font.size = Pt(18)
    cp2.font.color.rgb = text_gray
    cp2.font.name = "Consolas"
    
    # 底部页码区域
    footer_box = slide2.shapes.add_textbox(Inches(11.5), Inches(7), Inches(1.5), Inches(0.4))
    ftf = footer_box.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "[02/10]"
    fp.font.size = Pt(12)
    fp.font.color.rgb = glitch_purple
    fp.font.name = "Consolas"
    fp.alignment = PP_ALIGN.RIGHT
    
    # 保存
    output_path = os.path.join(TEMPLATES_DIR, "hackathon_neon.pptx")
    prs.save(output_path)
    print(f"✓ 创建模板: {output_path}")
    return output_path


def create_academic_latex_template():
    """创建学术LaTeX极简风格模板 - 纯白背景+黑色+学术蓝"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 颜色配置
    bg_white = RGBColor(255, 255, 255)    # 纯白背景
    text_black = RGBColor(0, 0, 0)        # 黑色
    academic_blue = RGBColor(0, 51, 102)  # 学术蓝 #003366
    line_gray = RGBColor(200, 200, 200)   # 浅灰分隔线
    
    # ========== 封面页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 纯白背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_white
    bg.line.fill.background()
    
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # 顶部细线
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(11.333), Pt(1)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = text_black
    top_line.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "论文标题"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = text_black
    p.font.name = "Times New Roman"
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题/作者信息
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.6))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.text = "Author Name | Institution | Date"
    sp.font.size = Pt(18)
    sp.font.color.rgb = academic_blue
    sp.font.name = "Times New Roman"
    sp.alignment = PP_ALIGN.CENTER
    
    # 底部细线
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(6),
        Inches(11.333), Pt(1)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = text_black
    bottom_line.line.fill.background()
    
    # ========== 内容页 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 纯白背景
    bg2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = bg_white
    bg2.line.fill.background()
    
    spTree2 = slide2.shapes._spTree
    sp2 = bg2._element
    spTree2.remove(sp2)
    spTree2.insert(2, sp2)
    
    # 顶部学术蓝细线
    header_line = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Pt(3)
    )
    header_line.fill.solid()
    header_line.fill.fore_color.rgb = academic_blue
    header_line.line.fill.background()
    
    # 页面标题
    title_box2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
    tf2 = title_box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "1. Introduction"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = text_black
    p2.font.name = "Times New Roman"
    
    # 标题下划线
    title_underline = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.2),
        Inches(3), Pt(1.5)
    )
    title_underline.fill.solid()
    title_underline.fill.fore_color.rgb = academic_blue
    title_underline.line.fill.background()
    
    # 内容区域提示
    content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = "• Bullet point 1\n• Bullet point 2\n• Bullet point 3"
    cp.font.size = Pt(20)
    cp.font.color.rgb = text_black
    cp.font.name = "Times New Roman"
    cp.line_spacing = 1.5
    
    # 底部页码
    footer_line = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.35),
        Inches(13.333), Pt(1)
    )
    footer_line.fill.solid()
    footer_line.fill.fore_color.rgb = line_gray
    footer_line.line.fill.background()
    
    page_num = slide2.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.4))
    pnf = page_num.text_frame
    pnp = pnf.paragraphs[0]
    pnp.text = "2"
    pnp.font.size = Pt(12)
    pnp.font.color.rgb = academic_blue
    pnp.font.name = "Times New Roman"
    pnp.alignment = PP_ALIGN.RIGHT
    
    # 保存
    output_path = os.path.join(TEMPLATES_DIR, "academic_latex.pptx")
    prs.save(output_path)
    print(f"✓ 创建模板: {output_path}")
    return output_path


def create_yc_pitch_deck_template():
    """创建YC Pitch Deck融资演讲风格模板 - 大标题+大数字+视觉焦点"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 颜色配置
    bg_white = RGBColor(255, 255, 255)    # 纯白背景
    primary_orange = RGBColor(255, 102, 0) # YC橙色
    text_black = RGBColor(33, 33, 33)      # 深黑
    text_gray = RGBColor(128, 128, 128)    # 灰色
    accent_blue = RGBColor(0, 122, 255)    # 强调蓝
    
    # ========== 封面页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 白色背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_white
    bg.line.fill.background()
    
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # 左上角橙色块（品牌标识区）
    brand_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.3), Inches(7.5)
    )
    brand_block.fill.solid()
    brand_block.fill.fore_color.rgb = primary_orange
    brand_block.line.fill.background()
    
    # 巨大的公司名/产品名
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "YOUR STARTUP"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = text_black
    p.font.name = "Arial Black"
    
    # 一句话描述
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
    ttf = tagline_box.text_frame
    tp = ttf.paragraphs[0]
    tp.text = "一句话说清楚你在做什么"
    tp.font.size = Pt(28)
    tp.font.color.rgb = text_gray
    tp.font.name = "Arial"
    
    # 底部信息
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
    ftf = footer_box.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "Seed Round · $2M · 2024"
    fp.font.size = Pt(16)
    fp.font.color.rgb = primary_orange
    fp.font.name = "Arial"
    
    # ========== 问题页 (Problem) ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = bg_white
    bg2.line.fill.background()
    
    spTree2 = slide2.shapes._spTree
    sp2 = bg2._element
    spTree2.remove(sp2)
    spTree2.insert(2, sp2)
    
    # 行动口号标题
    headline = slide2.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(1.5))
    htf = headline.text_frame
    htf.word_wrap = True
    hp = htf.paragraphs[0]
    hp.text = "THE PROBLEM"
    hp.font.size = Pt(48)
    hp.font.bold = True
    hp.font.color.rgb = text_black
    hp.font.name = "Arial Black"
    
    # 大数字视觉焦点
    big_number = slide2.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(6), Inches(2.5))
    bnf = big_number.text_frame
    bnp = bnf.paragraphs[0]
    bnp.text = "73%"
    bnp.font.size = Pt(144)
    bnp.font.bold = True
    bnp.font.color.rgb = primary_orange
    bnp.font.name = "Arial Black"
    
    # 数字说明
    number_desc = slide2.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(10), Inches(1))
    ndf = number_desc.text_frame
    ndp = ndf.paragraphs[0]
    ndp.text = "的用户面临这个痛点"
    ndp.font.size = Pt(24)
    ndp.font.color.rgb = text_gray
    ndp.font.name = "Arial"
    
    # 右侧装饰条
    right_bar = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(13.033), Inches(0),
        Inches(0.3), Inches(7.5)
    )
    right_bar.fill.solid()
    right_bar.fill.fore_color.rgb = primary_orange
    right_bar.line.fill.background()
    
    # 页码
    page_num = slide2.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.4))
    pnf = page_num.text_frame
    pnp = pnf.paragraphs[0]
    pnp.text = "02"
    pnp.font.size = Pt(14)
    pnp.font.color.rgb = text_gray
    pnp.font.name = "Arial"
    pnp.alignment = PP_ALIGN.RIGHT
    
    # ========== 解决方案页 (Solution) ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg3 = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = primary_orange  # 橙色背景突出
    bg3.line.fill.background()
    
    spTree3 = slide3.shapes._spTree
    sp3 = bg3._element
    spTree3.remove(sp3)
    spTree3.insert(2, sp3)
    
    # 白色大标题
    solution_title = slide3.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2))
    stf3 = solution_title.text_frame
    stf3.word_wrap = True
    sp3 = stf3.paragraphs[0]
    sp3.text = "WE FIX THIS."
    sp3.font.size = Pt(72)
    sp3.font.bold = True
    sp3.font.color.rgb = bg_white
    sp3.font.name = "Arial Black"
    
    # 副标题
    solution_sub = slide3.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.5), Inches(1))
    ssf = solution_sub.text_frame
    ssp = ssf.paragraphs[0]
    ssp.text = "用一句话描述你的解决方案"
    ssp.font.size = Pt(24)
    ssp.font.color.rgb = bg_white
    ssp.font.name = "Arial"
    
    # 保存
    output_path = os.path.join(TEMPLATES_DIR, "yc_pitch_deck.pptx")
    prs.save(output_path)
    print(f"✓ 创建模板: {output_path}")
    return output_path


def create_data_dashboard_template():
    """创建数据仪表盘风格模板 - Bento Box布局+深蓝灰背景+高对比图表色"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 颜色配置
    bg_dark = RGBColor(15, 23, 42)         # 深蓝灰 #0F172A
    card_bg = RGBColor(30, 41, 59)         # 卡片背景
    accent_yellow = RGBColor(250, 204, 21) # 亮黄 #FACC15
    accent_cyan = RGBColor(34, 211, 238)   # 青色 #22D3EE
    accent_green = RGBColor(74, 222, 128)  # 绿色
    accent_red = RGBColor(248, 113, 113)   # 红色
    text_white = RGBColor(255, 255, 255)
    text_gray = RGBColor(148, 163, 184)    # 浅灰文字
    
    # ========== 封面页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 深蓝灰背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_dark
    bg.line.fill.background()
    
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # 顶部渐变装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.08)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent_cyan
    top_bar.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DATA DASHBOARD"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = text_white
    p.font.name = "Arial"
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.8))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.text = "数据驱动决策 · 可视化洞察"
    sp.font.size = Pt(24)
    sp.font.color.rgb = accent_yellow
    sp.font.name = "Arial"
    
    # 底部装饰 - 模拟数据点
    for i, color in enumerate([accent_cyan, accent_yellow, accent_green, accent_red]):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8 + i * 0.6), Inches(6.5),
            Inches(0.15), Inches(0.15)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
    
    # ========== 仪表盘页 - Bento Box 布局 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 深蓝灰背景
    bg2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = bg_dark
    bg2.line.fill.background()
    
    spTree2 = slide2.shapes._spTree
    sp2 = bg2._element
    spTree2.remove(sp2)
    spTree2.insert(2, sp2)
    
    # 页面标题
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.6))
    tf2 = title2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Overview Dashboard"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = text_white
    p2.font.name = "Arial"
    
    # KPI 卡片 1 - 左上
    kpi1 = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.1),
        Inches(3), Inches(1.8)
    )
    kpi1.fill.solid()
    kpi1.fill.fore_color.rgb = card_bg
    kpi1.line.fill.background()
    
    kpi1_label = slide2.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(2.6), Inches(0.4))
    kpi1_ltf = kpi1_label.text_frame
    kpi1_lp = kpi1_ltf.paragraphs[0]
    kpi1_lp.text = "总访问量"
    kpi1_lp.font.size = Pt(14)
    kpi1_lp.font.color.rgb = text_gray
    
    kpi1_value = slide2.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(2.6), Inches(0.8))
    kpi1_vtf = kpi1_value.text_frame
    kpi1_vp = kpi1_vtf.paragraphs[0]
    kpi1_vp.text = "1.2M"
    kpi1_vp.font.size = Pt(42)
    kpi1_vp.font.bold = True
    kpi1_vp.font.color.rgb = accent_cyan
    
    # KPI 卡片 2 - 中上
    kpi2 = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.7), Inches(1.1),
        Inches(3), Inches(1.8)
    )
    kpi2.fill.solid()
    kpi2.fill.fore_color.rgb = card_bg
    kpi2.line.fill.background()
    
    kpi2_label = slide2.shapes.add_textbox(Inches(3.9), Inches(1.3), Inches(2.6), Inches(0.4))
    kpi2_ltf = kpi2_label.text_frame
    kpi2_lp = kpi2_ltf.paragraphs[0]
    kpi2_lp.text = "转化率"
    kpi2_lp.font.size = Pt(14)
    kpi2_lp.font.color.rgb = text_gray
    
    kpi2_value = slide2.shapes.add_textbox(Inches(3.9), Inches(1.7), Inches(2.6), Inches(0.8))
    kpi2_vtf = kpi2_value.text_frame
    kpi2_vp = kpi2_vtf.paragraphs[0]
    kpi2_vp.text = "23.5%"
    kpi2_vp.font.size = Pt(42)
    kpi2_vp.font.bold = True
    kpi2_vp.font.color.rgb = accent_yellow
    
    # KPI 卡片 3 - 右上
    kpi3 = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.9), Inches(1.1),
        Inches(3), Inches(1.8)
    )
    kpi3.fill.solid()
    kpi3.fill.fore_color.rgb = card_bg
    kpi3.line.fill.background()
    
    kpi3_label = slide2.shapes.add_textbox(Inches(7.1), Inches(1.3), Inches(2.6), Inches(0.4))
    kpi3_ltf = kpi3_label.text_frame
    kpi3_lp = kpi3_ltf.paragraphs[0]
    kpi3_lp.text = "收入增长"
    kpi3_lp.font.size = Pt(14)
    kpi3_lp.font.color.rgb = text_gray
    
    kpi3_value = slide2.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(2.6), Inches(0.8))
    kpi3_vtf = kpi3_value.text_frame
    kpi3_vp = kpi3_vtf.paragraphs[0]
    kpi3_vp.text = "+18.2%"
    kpi3_vp.font.size = Pt(42)
    kpi3_vp.font.bold = True
    kpi3_vp.font.color.rgb = accent_green
    
    # KPI 卡片 4 - 最右上
    kpi4 = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.1), Inches(1.1),
        Inches(2.7), Inches(1.8)
    )
    kpi4.fill.solid()
    kpi4.fill.fore_color.rgb = card_bg
    kpi4.line.fill.background()
    
    kpi4_label = slide2.shapes.add_textbox(Inches(10.3), Inches(1.3), Inches(2.3), Inches(0.4))
    kpi4_ltf = kpi4_label.text_frame
    kpi4_lp = kpi4_ltf.paragraphs[0]
    kpi4_lp.text = "跳出率"
    kpi4_lp.font.size = Pt(14)
    kpi4_lp.font.color.rgb = text_gray
    
    kpi4_value = slide2.shapes.add_textbox(Inches(10.3), Inches(1.7), Inches(2.3), Inches(0.8))
    kpi4_vtf = kpi4_value.text_frame
    kpi4_vp = kpi4_vtf.paragraphs[0]
    kpi4_vp.text = "32%"
    kpi4_vp.font.size = Pt(42)
    kpi4_vp.font.bold = True
    kpi4_vp.font.color.rgb = accent_red
    
    # 大图表区域 - 左下
    chart_area = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(3.1),
        Inches(8), Inches(4)
    )
    chart_area.fill.solid()
    chart_area.fill.fore_color.rgb = card_bg
    chart_area.line.fill.background()
    
    chart_title = slide2.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(4), Inches(0.4))
    ctf = chart_title.text_frame
    ctp = ctf.paragraphs[0]
    ctp.text = "趋势分析 - 折线图区域"
    ctp.font.size = Pt(14)
    ctp.font.color.rgb = text_gray
    
    # 模拟图表线条
    for i in range(5):
        bar = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.2 + i * 1.4), Inches(6.5 - (i % 3 + 1) * 0.8),
            Inches(0.8), Inches((i % 3 + 1) * 0.8)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_cyan if i % 2 == 0 else accent_yellow
        bar.line.fill.background()
    
    # 右侧小图表区域
    side_chart = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.7), Inches(3.1),
        Inches(4.1), Inches(4)
    )
    side_chart.fill.solid()
    side_chart.fill.fore_color.rgb = card_bg
    side_chart.line.fill.background()
    
    side_title = slide2.shapes.add_textbox(Inches(8.9), Inches(3.3), Inches(3.7), Inches(0.4))
    stf2 = side_title.text_frame
    stp = stf2.paragraphs[0]
    stp.text = "分布占比 - 饼图区域"
    stp.font.size = Pt(14)
    stp.font.color.rgb = text_gray
    
    # 模拟饼图
    pie = slide2.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(9.5), Inches(4.2),
        Inches(2.5), Inches(2.5)
    )
    pie.fill.solid()
    pie.fill.fore_color.rgb = accent_cyan
    pie.line.color.rgb = bg_dark
    pie.line.width = Pt(3)
    
    # 保存
    output_path = os.path.join(TEMPLATES_DIR, "data_dashboard.pptx")
    prs.save(output_path)
    print(f"✓ 创建模板: {output_path}")
    return output_path


def create_swiss_style_template():
    """创建瑞士国际主义设计风格模板 - 网格对齐+强烈字号对比+克莱因蓝/国际橙"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 颜色配置
    bg_white = RGBColor(255, 255, 255)     # 纯白背景
    klein_blue = RGBColor(0, 47, 167)      # 克莱因蓝 #002FA7
    intl_orange = RGBColor(255, 79, 0)     # 国际橙 #FF4F00
    text_black = RGBColor(0, 0, 0)         # 纯黑
    text_gray = RGBColor(100, 100, 100)    # 灰色
    
    # ========== 封面页 - 非对称布局 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 白色背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_white
    bg.line.fill.background()
    
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # 左侧大色块 - 克莱因蓝
    left_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(5), Inches(7.5)
    )
    left_block.fill.solid()
    left_block.fill.fore_color.rgb = klein_blue
    left_block.line.fill.background()
    
    # 超大标题数字/字母 - 在蓝色块上
    big_letter = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4), Inches(4))
    bltf = big_letter.text_frame
    blp = bltf.paragraphs[0]
    blp.text = "01"
    blp.font.size = Pt(200)
    blp.font.bold = True
    blp.font.color.rgb = bg_white
    blp.font.name = "Arial"
    
    # 右侧标题 - 极大字号
    title_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(7.5), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SWISS\nSTYLE"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = text_black
    p.font.name = "Arial"
    
    # 右下角小字 - 极小字号形成对比
    small_text = slide.shapes.add_textbox(Inches(5.5), Inches(6), Inches(7), Inches(1))
    stf = small_text.text_frame
    stp = stf.paragraphs[0]
    stp.text = "International Typographic Style\nGrid · Helvetica · Asymmetry"
    stp.font.size = Pt(12)
    stp.font.color.rgb = text_gray
    stp.font.name = "Arial"
    
    # ========== 内容页 - 网格布局 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 白色背景
    bg2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = bg_white
    bg2.line.fill.background()
    
    spTree2 = slide2.shapes._spTree
    sp2 = bg2._element
    spTree2.remove(sp2)
    spTree2.insert(2, sp2)
    
    # 顶部国际橙色条
    top_bar = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.5)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = intl_orange
    top_bar.line.fill.background()
    
    # 左侧竖排大标题
    vertical_title = slide2.shapes.add_textbox(Inches(0.3), Inches(1), Inches(1.5), Inches(6))
    vtf = vertical_title.text_frame
    vtp = vtf.paragraphs[0]
    vtp.text = "02"
    vtp.font.size = Pt(120)
    vtp.font.bold = True
    vtp.font.color.rgb = text_black
    vtp.font.name = "Arial"
    
    # 章节标题 - 右上
    section_title = slide2.shapes.add_textbox(Inches(2.5), Inches(1.2), Inches(10), Inches(1))
    sctf = section_title.text_frame
    sctp = sctf.paragraphs[0]
    sctp.text = "CHAPTER TITLE"
    sctp.font.size = Pt(48)
    sctp.font.bold = True
    sctp.font.color.rgb = text_black
    sctp.font.name = "Arial"
    
    # 网格参考线（视觉元素）
    for i in range(4):
        grid_line = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2.5 + i * 2.7), Inches(2.5),
            Pt(0.5), Inches(4.5)
        )
        grid_line.fill.solid()
        grid_line.fill.fore_color.rgb = RGBColor(230, 230, 230)
        grid_line.line.fill.background()
    
    # 内容区域 - 右下角，留白设计
    content_box = slide2.shapes.add_textbox(Inches(2.5), Inches(3), Inches(5), Inches(3.5))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    ctp = ctf.paragraphs[0]
    ctp.text = "Content goes here.\nStrict grid alignment.\nAsymmetric composition."
    ctp.font.size = Pt(18)
    ctp.font.color.rgb = text_black
    ctp.font.name = "Arial"
    ctp.line_spacing = 1.8
    
    # 右侧装饰色块
    right_accent = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(10), Inches(3),
        Inches(2.5), Inches(2.5)
    )
    right_accent.fill.solid()
    right_accent.fill.fore_color.rgb = klein_blue
    right_accent.line.fill.background()
    
    # 页码 - 底部右侧
    page_num = slide2.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.4))
    pnf = page_num.text_frame
    pnp = pnf.paragraphs[0]
    pnp.text = "02"
    pnp.font.size = Pt(10)
    pnp.font.color.rgb = text_gray
    pnp.font.name = "Arial"
    pnp.alignment = PP_ALIGN.RIGHT
    
    # 保存
    output_path = os.path.join(TEMPLATES_DIR, "swiss_style.pptx")
    prs.save(output_path)
    print(f"✓ 创建模板: {output_path}")
    return output_path


def create_futuristic_glass_template():
    """创建未来科技毛玻璃极光风格模板 - Glassmorphism+极光渐变+3D元素"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 颜色配置 - 极光色系
    aurora_purple = RGBColor(139, 92, 246)   # 紫色 #8B5CF6
    aurora_blue = RGBColor(59, 130, 246)     # 蓝色 #3B82F6
    aurora_cyan = RGBColor(34, 211, 238)     # 青色 #22D3EE
    aurora_pink = RGBColor(236, 72, 153)     # 粉色 #EC4899
    bg_dark = RGBColor(15, 10, 40)           # 深紫黑背景
    glass_bg = RGBColor(255, 255, 255)       # 毛玻璃白（会设置透明度）
    text_white = RGBColor(255, 255, 255)
    text_light = RGBColor(200, 200, 220)
    
    # ========== 封面页 - 极光背景 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 深色背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_dark
    bg.line.fill.background()
    
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # 极光效果 - 多个渐变色块模拟
    aurora1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-2), Inches(-1),
        Inches(8), Inches(5)
    )
    aurora1.fill.solid()
    aurora1.fill.fore_color.rgb = aurora_purple
    aurora1.line.fill.background()
    
    aurora2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(6), Inches(3),
        Inches(10), Inches(6)
    )
    aurora2.fill.solid()
    aurora2.fill.fore_color.rgb = aurora_blue
    aurora2.line.fill.background()
    
    aurora3 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(8), Inches(-2),
        Inches(6), Inches(4)
    )
    aurora3.fill.solid()
    aurora3.fill.fore_color.rgb = aurora_cyan
    aurora3.line.fill.background()
    
    # 3D 发光球体效果
    glow_sphere = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(9), Inches(1.5),
        Inches(3), Inches(3)
    )
    glow_sphere.fill.solid()
    glow_sphere.fill.fore_color.rgb = aurora_pink
    glow_sphere.line.fill.background()
    
    # 毛玻璃卡片 - 标题区域
    glass_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2),
        Inches(7), Inches(3.5)
    )
    glass_card.fill.solid()
    glass_card.fill.fore_color.rgb = RGBColor(40, 40, 60)
    glass_card.line.color.rgb = RGBColor(100, 100, 140)
    glass_card.line.width = Pt(1)
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(6), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "FUTURE\nINTERFACE"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = text_white
    p.font.name = "Arial"
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(6), Inches(0.8))
    stf = sub_box.text_frame
    stp = stf.paragraphs[0]
    stp.text = "Glassmorphism · Aurora · 3D Elements"
    stp.font.size = Pt(16)
    stp.font.color.rgb = aurora_cyan
    stp.font.name = "Arial"
    
    # 底部装饰线
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(6.8),
        Inches(4), Pt(2)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = aurora_cyan
    bottom_line.line.fill.background()
    
    # ========== 内容页 - 毛玻璃卡片布局 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 深色背景
    bg2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = bg_dark
    bg2.line.fill.background()
    
    spTree2 = slide2.shapes._spTree
    sp2 = bg2._element
    spTree2.remove(sp2)
    spTree2.insert(2, sp2)
    
    # 背景极光
    bg_aurora1 = slide2.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-3), Inches(4),
        Inches(8), Inches(6)
    )
    bg_aurora1.fill.solid()
    bg_aurora1.fill.fore_color.rgb = aurora_purple
    bg_aurora1.line.fill.background()
    
    bg_aurora2 = slide2.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(10), Inches(-2),
        Inches(6), Inches(5)
    )
    bg_aurora2.fill.solid()
    bg_aurora2.fill.fore_color.rgb = aurora_blue
    bg_aurora2.line.fill.background()
    
    # 页面标题
    title2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.8))
    tf2 = title2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "// CHAPTER_01"
    p2.font.size = Pt(14)
    p2.font.color.rgb = aurora_cyan
    p2.font.name = "Arial"
    
    section_title = slide2.shapes.add_textbox(Inches(0.8), Inches(1), Inches(10), Inches(1))
    sctf = section_title.text_frame
    sctp = sctf.paragraphs[0]
    sctp.text = "Section Title"
    sctp.font.size = Pt(36)
    sctp.font.bold = True
    sctp.font.color.rgb = text_white
    sctp.font.name = "Arial"
    
    # 毛玻璃内容卡片 1
    card1 = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(2.2),
        Inches(5.5), Inches(4.5)
    )
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(40, 40, 60)
    card1.line.color.rgb = RGBColor(100, 100, 140)
    card1.line.width = Pt(1)
    
    card1_title = slide2.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(4.8), Inches(0.6))
    c1tf = card1_title.text_frame
    c1p = c1tf.paragraphs[0]
    c1p.text = "Feature One"
    c1p.font.size = Pt(24)
    c1p.font.bold = True
    c1p.font.color.rgb = text_white
    
    card1_content = slide2.shapes.add_textbox(Inches(1.2), Inches(3.3), Inches(4.8), Inches(3))
    c1ctf = card1_content.text_frame
    c1ctf.word_wrap = True
    c1cp = c1ctf.paragraphs[0]
    c1cp.text = "Content description goes here.\nFuturistic and minimal."
    c1cp.font.size = Pt(14)
    c1cp.font.color.rgb = text_light
    
    # 毛玻璃内容卡片 2
    card2 = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(2.2),
        Inches(5.5), Inches(4.5)
    )
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(40, 40, 60)
    card2.line.color.rgb = RGBColor(100, 100, 140)
    card2.line.width = Pt(1)
    
    card2_title = slide2.shapes.add_textbox(Inches(7.2), Inches(2.6), Inches(4.8), Inches(0.6))
    c2tf = card2_title.text_frame
    c2p = c2tf.paragraphs[0]
    c2p.text = "Feature Two"
    c2p.font.size = Pt(24)
    c2p.font.bold = True
    c2p.font.color.rgb = text_white
    
    # 3D 装饰球体
    deco_sphere = slide2.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(11.5), Inches(5.5),
        Inches(1.5), Inches(1.5)
    )
    deco_sphere.fill.solid()
    deco_sphere.fill.fore_color.rgb = aurora_pink
    deco_sphere.line.fill.background()
    
    # 保存
    output_path = os.path.join(TEMPLATES_DIR, "futuristic_glass.pptx")
    prs.save(output_path)
    print(f"✓ 创建模板: {output_path}")
    return output_path


def create_retro_pixel_template():
    """创建8-bit像素复古游戏风格模板 - 像素风+街机字体+鲜艳原色"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 颜色配置 - 复古游戏原色
    pixel_red = RGBColor(255, 0, 77)       # 红色
    pixel_blue = RGBColor(41, 173, 255)    # 蓝色
    pixel_yellow = RGBColor(255, 236, 39)  # 黄色
    pixel_green = RGBColor(0, 228, 54)     # 绿色
    bg_dark = RGBColor(20, 12, 28)         # 深紫黑背景
    bg_blue = RGBColor(29, 43, 83)         # 深蓝背景
    text_white = RGBColor(255, 241, 232)   # 米白色
    
    # ========== 封面页 - 游戏开始画面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 深色背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_dark
    bg.line.fill.background()
    
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # 像素边框 - 顶部
    for i in range(27):
        pixel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(i * 0.5), Inches(0),
            Inches(0.4), Inches(0.4)
        )
        pixel.fill.solid()
        pixel.fill.fore_color.rgb = pixel_yellow if i % 2 == 0 else pixel_red
        pixel.line.fill.background()
    
    # 像素边框 - 底部
    for i in range(27):
        pixel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(i * 0.5), Inches(7.1),
            Inches(0.4), Inches(0.4)
        )
        pixel.fill.solid()
        pixel.fill.fore_color.rgb = pixel_blue if i % 2 == 0 else pixel_green
        pixel.line.fill.background()
    
    # 主标题 - 游戏风格
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PRESS START"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = pixel_yellow
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    stf = sub_box.text_frame
    stp = stf.paragraphs[0]
    stp.text = "YOUR GAME TITLE HERE"
    stp.font.size = Pt(28)
    stp.font.color.rgb = text_white
    stp.font.name = "Consolas"
    stp.alignment = PP_ALIGN.CENTER
    
    # 闪烁提示文字
    blink_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11.333), Inches(0.6))
    btf = blink_box.text_frame
    bp = btf.paragraphs[0]
    bp.text = "► INSERT COIN TO CONTINUE ◄"
    bp.font.size = Pt(18)
    bp.font.color.rgb = pixel_red
    bp.font.name = "Consolas"
    bp.alignment = PP_ALIGN.CENTER
    
    # 像素星星装饰
    star_positions = [(2, 1.5), (11, 1.2), (3, 5), (10, 5.5), (1, 3)]
    for x, y in star_positions:
        star = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(0.15), Inches(0.15)
        )
        star.fill.solid()
        star.fill.fore_color.rgb = text_white
        star.line.fill.background()
    
    # ========== 关卡选择页 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 深蓝背景
    bg2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = bg_blue
    bg2.line.fill.background()
    
    spTree2 = slide2.shapes._spTree
    sp2 = bg2._element
    spTree2.remove(sp2)
    spTree2.insert(2, sp2)
    
    # 页面标题
    level_title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    ltf = level_title.text_frame
    ltp = ltf.paragraphs[0]
    ltp.text = "★ LEVEL SELECT ★"
    ltp.font.size = Pt(32)
    ltp.font.bold = True
    ltp.font.color.rgb = pixel_yellow
    ltp.font.name = "Consolas"
    
    # 关卡方块 - 模拟关卡选择地图
    level_colors = [pixel_green, pixel_blue, pixel_yellow, pixel_red]
    level_names = ["LV.1", "LV.2", "LV.3", "LV.4"]
    for i in range(4):
        # 关卡方块
        level_block = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.5 + i * 2.8), Inches(2),
            Inches(2), Inches(2)
        )
        level_block.fill.solid()
        level_block.fill.fore_color.rgb = level_colors[i]
        level_block.line.color.rgb = text_white
        level_block.line.width = Pt(3)
        
        # 关卡编号
        level_num = slide2.shapes.add_textbox(
            Inches(1.5 + i * 2.8), Inches(2.6),
            Inches(2), Inches(0.8)
        )
        lnf = level_num.text_frame
        lnp = lnf.paragraphs[0]
        lnp.text = level_names[i]
        lnp.font.size = Pt(28)
        lnp.font.bold = True
        lnp.font.color.rgb = bg_dark
        lnp.font.name = "Consolas"
        lnp.alignment = PP_ALIGN.CENTER
    
    # 连接线（像素风格路径）
    for i in range(3):
        path = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3.5 + i * 2.8), Inches(2.9),
            Inches(0.8), Inches(0.2)
        )
        path.fill.solid()
        path.fill.fore_color.rgb = pixel_yellow
        path.line.fill.background()
    
    # 内容区域
    content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.5), Inches(2))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    ctp = ctf.paragraphs[0]
    ctp.text = "► SELECT YOUR LEVEL\n► EACH LEVEL = ONE CHAPTER\n► COLLECT ALL KNOWLEDGE COINS!"
    ctp.font.size = Pt(18)
    ctp.font.color.rgb = text_white
    ctp.font.name = "Consolas"
    ctp.line_spacing = 1.5
    
    # 底部状态栏
    status_bar = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7),
        Inches(13.333), Inches(0.5)
    )
    status_bar.fill.solid()
    status_bar.fill.fore_color.rgb = bg_dark
    status_bar.line.fill.background()
    
    # 生命值显示
    lives_box = slide2.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(4), Inches(0.4))
    lvf = lives_box.text_frame
    lvp = lvf.paragraphs[0]
    lvp.text = "♥♥♥ LIVES: 3"
    lvp.font.size = Pt(14)
    lvp.font.color.rgb = pixel_red
    lvp.font.name = "Consolas"
    
    # 分数显示
    score_box = slide2.shapes.add_textbox(Inches(9), Inches(7.05), Inches(4), Inches(0.4))
    scf = score_box.text_frame
    scp = scf.paragraphs[0]
    scp.text = "SCORE: 00000"
    scp.font.size = Pt(14)
    scp.font.color.rgb = pixel_yellow
    scp.font.name = "Consolas"
    scp.alignment = PP_ALIGN.RIGHT
    
    # 保存
    output_path = os.path.join(TEMPLATES_DIR, "retro_pixel.pptx")
    prs.save(output_path)
    print(f"✓ 创建模板: {output_path}")
    return output_path


def create_handdrawn_sketch_template():
    """创建手绘视觉笔记涂鸦风格模板 - 牛皮纸背景+手写体+草图风格"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 颜色配置 - 手绘风格
    paper_bg = RGBColor(245, 235, 220)     # 牛皮纸/米黄色
    ink_black = RGBColor(45, 45, 45)       # 墨水黑
    ink_blue = RGBColor(30, 100, 180)      # 钢笔蓝
    highlight_yellow = RGBColor(255, 230, 100)  # 荧光笔黄
    highlight_pink = RGBColor(255, 180, 180)    # 荧光笔粉
    pencil_gray = RGBColor(120, 120, 120)  # 铅笔灰
    
    # ========== 封面页 - 笔记本封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 牛皮纸背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = paper_bg
    bg.line.fill.background()
    
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # 笔记本横线
    for i in range(1, 14):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(i * 0.55),
            Inches(12.333), Pt(0.5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(200, 190, 175)
        line.line.fill.background()
    
    # 左侧红色边线（笔记本风格）
    red_margin = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.2), Inches(0),
        Pt(2), Inches(7.5)
    )
    red_margin.fill.solid()
    red_margin.fill.fore_color.rgb = RGBColor(220, 100, 100)
    red_margin.line.fill.background()
    
    # 手绘风格标题框（不规则矩形模拟）
    title_box_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(2),
        Inches(9), Inches(2.5)
    )
    title_box_bg.fill.solid()
    title_box_bg.fill.fore_color.rgb = highlight_yellow
    title_box_bg.line.color.rgb = ink_black
    title_box_bg.line.width = Pt(2)
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(2.3), Inches(2.3), Inches(8.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Visual Notes ✏️"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = ink_black
    p.font.name = "Comic Sans MS"
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(2.3), Inches(3.8), Inches(8.5), Inches(0.8))
    stf = sub_box.text_frame
    stp = stf.paragraphs[0]
    stp.text = "~ 手绘笔记 · 轻松学习 ~"
    stp.font.size = Pt(24)
    stp.font.color.rgb = ink_blue
    stp.font.name = "Comic Sans MS"
    stp.alignment = PP_ALIGN.CENTER
    
    # 涂鸦装饰 - 灯泡（代表灵感）
    bulb = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(10.5), Inches(1),
        Inches(1.2), Inches(1.2)
    )
    bulb.fill.solid()
    bulb.fill.fore_color.rgb = highlight_yellow
    bulb.line.color.rgb = ink_black
    bulb.line.width = Pt(2)
    
    # 灯泡底座
    bulb_base = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(10.7), Inches(2.1),
        Inches(0.8), Inches(0.3)
    )
    bulb_base.fill.solid()
    bulb_base.fill.fore_color.rgb = pencil_gray
    bulb_base.line.color.rgb = ink_black
    bulb_base.line.width = Pt(1)
    
    # 涂鸦星星
    star_positions = [(1.5, 1), (11.5, 5.5), (0.8, 5)]
    for x, y in star_positions:
        star = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.5), Inches(0.5))
        stf = star.text_frame
        sp = stf.paragraphs[0]
        sp.text = "★"
        sp.font.size = Pt(24)
        sp.font.color.rgb = ink_blue
    
    # 底部日期
    date_box = slide.shapes.add_textbox(Inches(9), Inches(6.5), Inches(3.5), Inches(0.5))
    dtf = date_box.text_frame
    dp = dtf.paragraphs[0]
    dp.text = "Date: ___/___/___"
    dp.font.size = Pt(14)
    dp.font.color.rgb = pencil_gray
    dp.font.name = "Comic Sans MS"
    
    # ========== 内容页 - 笔记页面 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 牛皮纸背景
    bg2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = paper_bg
    bg2.line.fill.background()
    
    spTree2 = slide2.shapes._spTree
    sp2 = bg2._element
    spTree2.remove(sp2)
    spTree2.insert(2, sp2)
    
    # 笔记本横线
    for i in range(1, 14):
        line = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(i * 0.55),
            Inches(12.333), Pt(0.5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(200, 190, 175)
        line.line.fill.background()
    
    # 左侧红色边线
    red_margin2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.2), Inches(0),
        Pt(2), Inches(7.5)
    )
    red_margin2.fill.solid()
    red_margin2.fill.fore_color.rgb = RGBColor(220, 100, 100)
    red_margin2.line.fill.background()
    
    # 页面标题 - 带下划线
    title2 = slide2.shapes.add_textbox(Inches(1.5), Inches(0.3), Inches(10), Inches(0.8))
    tf2 = title2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Chapter 1: 重点笔记 📝"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = ink_black
    p2.font.name = "Comic Sans MS"
    
    # 手绘下划线
    underline = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(1),
        Inches(6), Pt(3)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = ink_blue
    underline.line.fill.background()
    
    # 荧光笔高亮区域
    highlight_box = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(1.5),
        Inches(5), Inches(1.5)
    )
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = highlight_yellow
    highlight_box.line.fill.background()
    
    # 重点内容
    key_point = slide2.shapes.add_textbox(Inches(1.6), Inches(1.6), Inches(4.8), Inches(1.3))
    kptf = key_point.text_frame
    kptf.word_wrap = True
    kpp = kptf.paragraphs[0]
    kpp.text = "💡 Key Point:\n这里写重点内容..."
    kpp.font.size = Pt(18)
    kpp.font.color.rgb = ink_black
    kpp.font.name = "Comic Sans MS"
    
    # 手绘箭头区域
    arrow_box = slide2.shapes.add_textbox(Inches(7), Inches(1.8), Inches(1), Inches(0.5))
    atf = arrow_box.text_frame
    ap = atf.paragraphs[0]
    ap.text = "→→→"
    ap.font.size = Pt(24)
    ap.font.color.rgb = ink_blue
    
    # 右侧涂鸦框
    doodle_box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8), Inches(1.3),
        Inches(4.5), Inches(2.5)
    )
    doodle_box.fill.solid()
    doodle_box.fill.fore_color.rgb = highlight_pink
    doodle_box.line.color.rgb = ink_black
    doodle_box.line.width = Pt(2)
    
    doodle_text = slide2.shapes.add_textbox(Inches(8.2), Inches(1.5), Inches(4.1), Inches(2.2))
    dtf2 = doodle_text.text_frame
    dtf2.word_wrap = True
    dp2 = dtf2.paragraphs[0]
    dp2.text = "🎨 Doodle Area\n\n[在这里画草图]"
    dp2.font.size = Pt(16)
    dp2.font.color.rgb = ink_black
    dp2.font.name = "Comic Sans MS"
    
    # 底部笔记列表
    notes_box = slide2.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(2.5))
    ntf = notes_box.text_frame
    ntf.word_wrap = True
    np = ntf.paragraphs[0]
    np.text = "✓ 要点一：_______________\n✓ 要点二：_______________\n✓ 要点三：_______________"
    np.font.size = Pt(18)
    np.font.color.rgb = ink_black
    np.font.name = "Comic Sans MS"
    np.line_spacing = 1.8
    
    # 页码
    page_num = slide2.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.4))
    pnf = page_num.text_frame
    pnp = pnf.paragraphs[0]
    pnp.text = "p.2"
    pnp.font.size = Pt(12)
    pnp.font.color.rgb = pencil_gray
    pnp.font.name = "Comic Sans MS"
    
    # 保存
    output_path = os.path.join(TEMPLATES_DIR, "handdrawn_sketch.pptx")
    prs.save(output_path)
    print(f"✓ 创建模板: {output_path}")
    return output_path


def create_mckinsey_consulting_template():
    """创建麦肯锡咨询风格模板 - 观点句标题+MECE结构+左文右图"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 颜色配置 - 专业商务
    bg_white = RGBColor(255, 255, 255)     # 纯白背景
    mckinsey_blue = RGBColor(0, 51, 141)   # 麦肯锡蓝
    text_black = RGBColor(51, 51, 51)      # 深灰黑
    text_gray = RGBColor(128, 128, 128)    # 灰色
    accent_gold = RGBColor(180, 151, 90)   # 金色强调
    line_gray = RGBColor(220, 220, 220)    # 浅灰线条
    
    # ========== 封面页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 白色背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_white
    bg.line.fill.background()
    
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # 顶部蓝色条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = mckinsey_blue
    top_bar.line.fill.background()
    
    # 左侧蓝色装饰块
    left_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.4), Inches(7.5)
    )
    left_block.fill.solid()
    left_block.fill.fore_color.rgb = mckinsey_blue
    left_block.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Strategic Report Title"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = text_black
    p.font.name = "Arial"
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(0.8))
    stf = sub_box.text_frame
    stp = stf.paragraphs[0]
    stp.text = "Subtitle | Client Name | Date"
    stp.font.size = Pt(20)
    stp.font.color.rgb = text_gray
    stp.font.name = "Arial"
    
    # 底部分隔线
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(5.5),
        Inches(4), Pt(2)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = accent_gold
    bottom_line.line.fill.background()
    
    # 机密标识
    conf_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(5), Inches(0.4))
    ctf = conf_box.text_frame
    cp = ctf.paragraphs[0]
    cp.text = "CONFIDENTIAL"
    cp.font.size = Pt(10)
    cp.font.color.rgb = text_gray
    cp.font.name = "Arial"
    
    # ========== 内容页 - 左文右图布局 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 白色背景
    bg2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = bg_white
    bg2.line.fill.background()
    
    spTree2 = slide2.shapes._spTree
    sp2 = bg2._element
    spTree2.remove(sp2)
    spTree2.insert(2, sp2)
    
    # 顶部蓝色条
    top_bar2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.08)
    )
    top_bar2.fill.solid()
    top_bar2.fill.fore_color.rgb = mckinsey_blue
    top_bar2.line.fill.background()
    
    # 观点句标题 (Action Title)
    action_title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    atf = action_title.text_frame
    atf.word_wrap = True
    atp = atf.paragraphs[0]
    atp.text = "观点句标题：完整表达本页核心观点的一句话"
    atp.font.size = Pt(24)
    atp.font.bold = True
    atp.font.color.rgb = mckinsey_blue
    atp.font.name = "Arial"
    
    # 标题下划线
    title_line = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.1),
        Inches(12.333), Pt(1)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = line_gray
    title_line.line.fill.background()
    
    # 左侧文字区域
    left_content = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
    lctf = left_content.text_frame
    lctf.word_wrap = True
    
    # 论据1
    lcp1 = lctf.paragraphs[0]
    lcp1.text = "■ 论据一"
    lcp1.font.size = Pt(16)
    lcp1.font.bold = True
    lcp1.font.color.rgb = text_black
    lcp1.font.name = "Arial"
    
    lcp1_detail = lctf.add_paragraph()
    lcp1_detail.text = "   支撑观点的具体数据或事实描述"
    lcp1_detail.font.size = Pt(14)
    lcp1_detail.font.color.rgb = text_gray
    lcp1_detail.font.name = "Arial"
    lcp1_detail.space_after = Pt(12)
    
    # 论据2
    lcp2 = lctf.add_paragraph()
    lcp2.text = "■ 论据二"
    lcp2.font.size = Pt(16)
    lcp2.font.bold = True
    lcp2.font.color.rgb = text_black
    lcp2.font.name = "Arial"
    
    lcp2_detail = lctf.add_paragraph()
    lcp2_detail.text = "   支撑观点的具体数据或事实描述"
    lcp2_detail.font.size = Pt(14)
    lcp2_detail.font.color.rgb = text_gray
    lcp2_detail.font.name = "Arial"
    lcp2_detail.space_after = Pt(12)
    
    # 论据3
    lcp3 = lctf.add_paragraph()
    lcp3.text = "■ 论据三"
    lcp3.font.size = Pt(16)
    lcp3.font.bold = True
    lcp3.font.color.rgb = text_black
    lcp3.font.name = "Arial"
    
    lcp3_detail = lctf.add_paragraph()
    lcp3_detail.text = "   支撑观点的具体数据或事实描述"
    lcp3_detail.font.size = Pt(14)
    lcp3_detail.font.color.rgb = text_gray
    lcp3_detail.font.name = "Arial"
    
    # 右侧图片占位区
    img_placeholder = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7), Inches(1.5),
        Inches(5.8), Inches(4.5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = RGBColor(245, 245, 245)
    img_placeholder.line.color.rgb = line_gray
    img_placeholder.line.width = Pt(1)
    
    img_text = slide2.shapes.add_textbox(Inches(7), Inches(3.5), Inches(5.8), Inches(0.8))
    itf = img_text.text_frame
    itp = itf.paragraphs[0]
    itp.text = "[图表/数据可视化区域]"
    itp.font.size = Pt(16)
    itp.font.color.rgb = text_gray
    itp.font.name = "Arial"
    itp.alignment = PP_ALIGN.CENTER
    
    # 底部分隔线
    footer_line = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(6.8),
        Inches(12.333), Pt(1)
    )
    footer_line.fill.solid()
    footer_line.fill.fore_color.rgb = line_gray
    footer_line.line.fill.background()
    
    # 页码
    page_num = slide2.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.4))
    pnf = page_num.text_frame
    pnp = pnf.paragraphs[0]
    pnp.text = "2"
    pnp.font.size = Pt(10)
    pnp.font.color.rgb = text_gray
    pnp.font.name = "Arial"
    pnp.alignment = PP_ALIGN.RIGHT
    
    # 来源标注
    source_box = slide2.shapes.add_textbox(Inches(0.5), Inches(7), Inches(8), Inches(0.4))
    srf = source_box.text_frame
    srp = srf.paragraphs[0]
    srp.text = "Source: [数据来源]"
    srp.font.size = Pt(8)
    srp.font.color.rgb = text_gray
    srp.font.name = "Arial"
    
    # 保存
    output_path = os.path.join(TEMPLATES_DIR, "mckinsey_consulting.pptx")
    prs.save(output_path)
    print(f"✓ 创建模板: {output_path}")
    return output_path


def create_magazine_editorial_template():
    """创建Vogue/Monocle杂志排版风格模板 - 大图+衬线体+高级叙事感"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 颜色配置 - 高级杂志风
    bg_white = RGBColor(255, 255, 255)     # 纯白
    bg_cream = RGBColor(250, 248, 245)     # 米白/奶油色
    text_black = RGBColor(30, 30, 30)      # 深黑
    text_gray = RGBColor(120, 120, 120)    # 灰色
    accent_gold = RGBColor(180, 160, 120)  # 金色
    img_placeholder = RGBColor(200, 195, 190)  # 图片占位灰
    
    # ========== 封面页 - 全屏图片+悬浮标题 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 图片占位区（全屏）
    full_img = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    full_img.fill.solid()
    full_img.fill.fore_color.rgb = img_placeholder
    full_img.line.fill.background()
    
    # 半透明遮罩（底部渐变效果模拟）
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(4.5),
        Inches(13.333), Inches(3)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(20, 20, 20)
    overlay.line.fill.background()
    
    # 杂志名/品牌标识
    brand_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(4), Inches(0.6))
    btf = brand_box.text_frame
    bp = btf.paragraphs[0]
    bp.text = "STORYTELLING"
    bp.font.size = Pt(14)
    bp.font.color.rgb = bg_white
    bp.font.name = "Georgia"
    bp.font.bold = True
    
    # 主标题 - 极大纤细衬线体
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Art of Visual Narrative"
    p.font.size = Pt(60)
    p.font.color.rgb = bg_white
    p.font.name = "Georgia"
    
    # 副标题/引言
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(8), Inches(0.8))
    stf = sub_box.text_frame
    stp = stf.paragraphs[0]
    stp.text = "A journey through images and words"
    stp.font.size = Pt(18)
    stp.font.italic = True
    stp.font.color.rgb = RGBColor(200, 200, 200)
    stp.font.name = "Georgia"
    
    # 期刊信息
    issue_box = slide.shapes.add_textbox(Inches(10.5), Inches(6.8), Inches(2.5), Inches(0.5))
    itf = issue_box.text_frame
    ip = itf.paragraphs[0]
    ip.text = "Issue 01 / 2025"
    ip.font.size = Pt(10)
    ip.font.color.rgb = RGBColor(180, 180, 180)
    ip.font.name = "Georgia"
    ip.alignment = PP_ALIGN.RIGHT
    
    # ========== 内容页 - 左图右文布局 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 奶油色背景
    bg2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = bg_cream
    bg2.line.fill.background()
    
    spTree2 = slide2.shapes._spTree
    sp2 = bg2._element
    spTree2.remove(sp2)
    spTree2.insert(2, sp2)
    
    # 左侧大图区域（占2/3）
    left_img = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(8.5), Inches(7.5)
    )
    left_img.fill.solid()
    left_img.fill.fore_color.rgb = img_placeholder
    left_img.line.fill.background()
    
    # 图片说明
    img_caption = slide2.shapes.add_textbox(Inches(0.5), Inches(7), Inches(7.5), Inches(0.4))
    ictf = img_caption.text_frame
    icp = ictf.paragraphs[0]
    icp.text = "[Full-bleed image placeholder]"
    icp.font.size = Pt(10)
    icp.font.italic = True
    icp.font.color.rgb = text_gray
    icp.font.name = "Georgia"
    
    # 右侧文字区域
    # 章节编号
    chapter_num = slide2.shapes.add_textbox(Inches(9), Inches(1), Inches(3.8), Inches(0.8))
    cntf = chapter_num.text_frame
    cnp = cntf.paragraphs[0]
    cnp.text = "01"
    cnp.font.size = Pt(72)
    cnp.font.color.rgb = accent_gold
    cnp.font.name = "Georgia"
    
    # 章节标题
    chapter_title = slide2.shapes.add_textbox(Inches(9), Inches(2.2), Inches(3.8), Inches(1.2))
    cttf = chapter_title.text_frame
    cttf.word_wrap = True
    ctp = cttf.paragraphs[0]
    ctp.text = "Chapter\nTitle"
    ctp.font.size = Pt(36)
    ctp.font.color.rgb = text_black
    ctp.font.name = "Georgia"
    
    # 分隔线
    divider = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(9), Inches(3.8),
        Inches(2), Pt(1)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = accent_gold
    divider.line.fill.background()
    
    # 正文内容
    body_text = slide2.shapes.add_textbox(Inches(9), Inches(4.2), Inches(3.8), Inches(2.5))
    bttf = body_text.text_frame
    bttf.word_wrap = True
    btp = bttf.paragraphs[0]
    btp.text = "Elegant body text goes here. Keep it minimal and let the image speak. Every word should earn its place on the page."
    btp.font.size = Pt(14)
    btp.font.color.rgb = text_gray
    btp.font.name = "Georgia"
    btp.line_spacing = 1.6
    
    # 页码
    page_num = slide2.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.4))
    pnf = page_num.text_frame
    pnp = pnf.paragraphs[0]
    pnp.text = "02"
    pnp.font.size = Pt(10)
    pnp.font.color.rgb = text_gray
    pnp.font.name = "Georgia"
    pnp.alignment = PP_ALIGN.RIGHT
    
    # ========== 内容页2 - 全屏图片+悬浮文字 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 全屏图片
    full_img3 = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    full_img3.fill.solid()
    full_img3.fill.fore_color.rgb = RGBColor(180, 175, 170)
    full_img3.line.fill.background()
    
    # 悬浮引言框（左上角）
    quote_bg = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.5),
        Inches(5), Inches(2.5)
    )
    quote_bg.fill.solid()
    quote_bg.fill.fore_color.rgb = bg_white
    quote_bg.line.fill.background()
    
    # 引言文字
    quote_text = slide3.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(4.4), Inches(2))
    qtf = quote_text.text_frame
    qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qp.text = '"A powerful quote that captures the essence of this moment."'
    qp.font.size = Pt(20)
    qp.font.italic = True
    qp.font.color.rgb = text_black
    qp.font.name = "Georgia"
    
    # 引言作者
    quote_author = slide3.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(4.4), Inches(0.4))
    qatf = quote_author.text_frame
    qap = qatf.paragraphs[0]
    qap.text = "— Author Name"
    qap.font.size = Pt(12)
    qap.font.color.rgb = accent_gold
    qap.font.name = "Georgia"
    
    # 页码
    page_num3 = slide3.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.4))
    pnf3 = page_num3.text_frame
    pnp3 = pnf3.paragraphs[0]
    pnp3.text = "03"
    pnp3.font.size = Pt(10)
    pnp3.font.color.rgb = bg_white
    pnp3.font.name = "Georgia"
    pnp3.alignment = PP_ALIGN.RIGHT
    
    # 保存
    output_path = os.path.join(TEMPLATES_DIR, "magazine_editorial.pptx")
    prs.save(output_path)
    print(f"✓ 创建模板: {output_path}")
    return output_path


def create_all_templates():
    """创建所有模板（仅亮色风格）"""
    os.makedirs(TEMPLATES_DIR, exist_ok=True)
    
    print("=" * 50)
    print("开始创建 PPT 模板...")
    print("=" * 50)
    
    create_academic_latex_template()
    create_yc_pitch_deck_template()
    create_swiss_style_template()
    create_handdrawn_sketch_template()
    create_mckinsey_consulting_template()
    create_magazine_editorial_template()
    
    print("=" * 50)
    print("✓ 所有模板创建完成！")
    print(f"模板目录: {TEMPLATES_DIR}")
    print("=" * 50)


if __name__ == "__main__":
    create_all_templates()
