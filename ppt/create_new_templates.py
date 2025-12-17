"""创建新模板 - 暗色主题、极简风格、中国风"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os


def create_dark_theme_template():
    """创建暗色主题模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 暗色配色方案
    bg_color = RGBColor(18, 18, 18)  # 深黑背景
    primary = RGBColor(96, 165, 250)  # 蓝色
    accent = RGBColor(244, 114, 182)  # 粉色
    text_primary = RGBColor(248, 250, 252)  # 白色文字
    text_secondary = RGBColor(148, 163, 184)  # 灰色文字

    # 创建封面布局
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 深色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # 渐变装饰
    accent_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(8), Inches(8)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = RGBColor(59, 130, 246)
    accent_shape.fill.fore_color.brightness = -0.3
    accent_shape.line.fill.background()

    prs.save('ppt/pptx_templates/dark_theme.pptx')
    print("✓ 暗色主题模板已创建: dark_theme.pptx")


def create_minimalist_template():
    """创建极简风格模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 极简配色
    bg_color = RGBColor(255, 255, 255)  # 纯白背景
    primary = RGBColor(17, 17, 17)  # 纯黑
    accent = RGBColor(229, 231, 235)  # 浅灰线条

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 白色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # 简洁底部线条
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), prs.slide_height - Inches(0.5),
        prs.slide_width - Inches(1), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = primary
    line.line.fill.background()

    prs.save('ppt/pptx_templates/minimalist.pptx')
    print("✓ 极简风格模板已创建: minimalist.pptx")


def create_chinese_style_template():
    """创建中国风模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 中国风配色
    bg_color = RGBColor(254, 252, 248)  # 米白色
    primary = RGBColor(180, 50, 50)  # 中国红
    gold = RGBColor(184, 134, 11)  # 金色
    ink = RGBColor(45, 45, 45)  # 墨色

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 米白背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # 红色边框装饰
    top_border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.3),
        prs.slide_width - Inches(0.6), Pt(4)
    )
    top_border.fill.solid()
    top_border.fill.fore_color.rgb = primary
    top_border.line.fill.background()

    bottom_border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.3), prs.slide_height - Inches(0.4),
        prs.slide_width - Inches(0.6), Pt(4)
    )
    bottom_border.fill.solid()
    bottom_border.fill.fore_color.rgb = primary
    bottom_border.line.fill.background()

    # 角落装饰
    corner_size = Inches(0.8)
    corners = [
        (Inches(0.3), Inches(0.3)),  # 左上
        (prs.slide_width - Inches(1.1), Inches(0.3)),  # 右上
        (Inches(0.3), prs.slide_height - Inches(1.1)),  # 左下
        (prs.slide_width - Inches(1.1), prs.slide_height - Inches(1.1)),  # 右下
    ]

    for x, y in corners:
        corner = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, corner_size, corner_size
        )
        corner.fill.solid()
        corner.fill.fore_color.rgb = gold
        corner.line.fill.background()

    prs.save('ppt/pptx_templates/chinese_style.pptx')
    print("✓ 中国风模板已创建: chinese_style.pptx")


def create_gradient_blue_template():
    """创建渐变蓝模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 渐变背景效果（使用多个矩形模拟）
    colors = [
        RGBColor(30, 64, 175),   # 深蓝
        RGBColor(59, 130, 246),  # 中蓝
        RGBColor(96, 165, 250),  # 浅蓝
    ]

    height_per_section = prs.slide_height / len(colors)
    for i, color in enumerate(colors):
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, height_per_section * i,
            prs.slide_width, height_per_section + Pt(2)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()

    prs.save('ppt/pptx_templates/gradient_blue.pptx')
    print("✓ 渐变蓝模板已创建: gradient_blue.pptx")


def create_tech_modern_template():
    """创建科技现代模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 深色科技背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42)
    bg.line.fill.background()

    # 网格线装饰（水平线）
    for i in range(1, 8):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(i),
            prs.slide_width, Pt(1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(51, 65, 85)
        line.line.fill.background()

    # 左侧发光条
    glow = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Pt(4), prs.slide_height
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(34, 211, 238)  # 青色
    glow.line.fill.background()

    prs.save('ppt/pptx_templates/tech_modern.pptx')
    print("✓ 科技现代模板已创建: tech_modern.pptx")


def create_warm_sunset_template():
    """创建暖色夕阳模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 暖色渐变
    colors = [
        RGBColor(251, 146, 60),   # 橙色
        RGBColor(249, 115, 22),   # 深橙
        RGBColor(234, 88, 12),    # 红橙
    ]

    height_per_section = prs.slide_height / len(colors)
    for i, color in enumerate(colors):
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, height_per_section * i,
            prs.slide_width, height_per_section + Pt(2)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()

    prs.save('ppt/pptx_templates/warm_sunset.pptx')
    print("✓ 暖色夕阳模板已创建: warm_sunset.pptx")


def create_corporate_blue_template():
    """创建企业蓝模板（正式商务）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 企业蓝配色
    bg_color = RGBColor(255, 255, 255)
    primary = RGBColor(0, 82, 147)  # 深蓝
    secondary = RGBColor(240, 245, 250)  # 浅蓝灰
    accent = RGBColor(0, 120, 200)  # 亮蓝

    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 白色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # 顶部深蓝色条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, Inches(1.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = primary
    top_bar.line.fill.background()

    # 底部细条
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.1),
        prs.slide_width, Inches(0.1)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = accent
    bottom_bar.line.fill.background()

    # 右侧装饰块
    right_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, prs.slide_width - Inches(1.5), Inches(1.2),
        Inches(1.5), prs.slide_height - Inches(1.3)
    )
    right_block.fill.solid()
    right_block.fill.fore_color.rgb = secondary
    right_block.line.fill.background()

    prs.save('ppt/pptx_templates/corporate_blue.pptx')
    print("✓ 企业蓝模板已创建: corporate_blue.pptx")


def create_nature_green_template():
    """创建自然绿模板（环保主题）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 自然绿配色
    bg_color = RGBColor(250, 255, 250)  # 淡绿白
    primary = RGBColor(34, 139, 34)  # 森林绿
    secondary = RGBColor(144, 238, 144)  # 浅绿
    accent = RGBColor(85, 170, 85)  # 草绿

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # 左侧渐变装饰
    left_accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(0.25), prs.slide_height
    )
    left_accent.fill.solid()
    left_accent.fill.fore_color.rgb = primary
    left_accent.line.fill.background()

    # 底部装饰条
    bottom_strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.8),
        prs.slide_width, Inches(0.8)
    )
    bottom_strip.fill.solid()
    bottom_strip.fill.fore_color.rgb = secondary
    bottom_strip.line.fill.background()

    # 圆形装饰
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, prs.slide_width - Inches(4), Inches(-1),
        Inches(5), Inches(5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(200, 240, 200)
    circle.line.fill.background()

    prs.save('ppt/pptx_templates/nature_green.pptx')
    print("✓ 自然绿模板已创建: nature_green.pptx")


def create_elegant_purple_template():
    """创建优雅紫模板（创意艺术）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 优雅紫配色
    primary = RGBColor(102, 51, 153)  # 紫色
    secondary = RGBColor(180, 150, 200)  # 淡紫
    bg_color = RGBColor(255, 250, 255)  # 淡紫白
    gold = RGBColor(200, 170, 110)  # 金色

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # 顶部紫色装饰
    top_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, Inches(0.15)
    )
    top_shape.fill.solid()
    top_shape.fill.fore_color.rgb = primary
    top_shape.line.fill.background()

    # 左侧装饰条
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.5),
        Pt(5), Inches(1.2)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = gold
    left_bar.line.fill.background()

    # 底部渐变装饰
    bottom_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(1),
        prs.slide_width, Inches(1)
    )
    bottom_shape.fill.solid()
    bottom_shape.fill.fore_color.rgb = secondary
    bottom_shape.line.fill.background()

    prs.save('ppt/pptx_templates/elegant_purple.pptx')
    print("✓ 优雅紫模板已创建: elegant_purple.pptx")


def create_startup_neon_template():
    """创建创业霓虹模板（现代活力）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 霓虹配色
    bg_color = RGBColor(20, 20, 30)  # 深色背景
    neon_pink = RGBColor(255, 0, 128)  # 霓虹粉
    neon_blue = RGBColor(0, 200, 255)  # 霓虹蓝
    neon_green = RGBColor(0, 255, 150)  # 霓虹绿

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 深色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # 霓虹粉线条
    pink_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2),
        prs.slide_width, Pt(3)
    )
    pink_line.fill.solid()
    pink_line.fill.fore_color.rgb = neon_pink
    pink_line.line.fill.background()

    # 霓虹蓝线条
    blue_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(5.5),
        prs.slide_width, Pt(3)
    )
    blue_line.fill.solid()
    blue_line.fill.fore_color.rgb = neon_blue
    blue_line.line.fill.background()

    # 左侧霓虹绿装饰
    green_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Pt(4), prs.slide_height
    )
    green_bar.fill.solid()
    green_bar.fill.fore_color.rgb = neon_green
    green_bar.line.fill.background()

    prs.save('ppt/pptx_templates/startup_neon.pptx')
    print("✓ 创业霓虹模板已创建: startup_neon.pptx")


def create_education_light_template():
    """创建教育亮色模板（学习主题）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 教育配色
    bg_color = RGBColor(255, 255, 255)
    primary = RGBColor(66, 133, 244)  # 蓝色
    secondary = RGBColor(251, 188, 4)  # 黄色
    accent = RGBColor(234, 67, 53)  # 红色
    green = RGBColor(52, 168, 83)  # 绿色

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 白色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # 顶部彩色条
    colors = [primary, secondary, accent, green]
    bar_width = prs.slide_width / 4
    for i, color in enumerate(colors):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, bar_width * i, 0,
            bar_width + Pt(1), Pt(8)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

    # 左侧蓝色装饰
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.5),
        Pt(6), Inches(1)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = primary
    left_bar.line.fill.background()

    # 底部装饰
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.3),
        prs.slide_width, Inches(0.3)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = RGBColor(240, 245, 250)
    bottom_bar.line.fill.background()

    prs.save('ppt/pptx_templates/education_light.pptx')
    print("✓ 教育亮色模板已创建: education_light.pptx")


def create_medical_clean_template():
    """创建医疗简洁模板（专业医疗）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 医疗配色
    bg_color = RGBColor(255, 255, 255)
    primary = RGBColor(0, 150, 136)  # 青绿色
    secondary = RGBColor(240, 250, 249)  # 淡青
    accent = RGBColor(0, 188, 212)  # 青色

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 白色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # 顶部青绿色条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, Inches(0.8)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = primary
    top_bar.line.fill.background()

    # 右侧装饰块
    right_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, prs.slide_width - Inches(2), Inches(0.8),
        Inches(2), prs.slide_height - Inches(0.8)
    )
    right_block.fill.solid()
    right_block.fill.fore_color.rgb = secondary
    right_block.line.fill.background()

    # 底部细线
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Pt(4),
        prs.slide_width, Pt(4)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = accent
    bottom_line.line.fill.background()

    # 十字装饰
    cross_h = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(11.5), Inches(0.2),
        Inches(0.6), Pt(5)
    )
    cross_h.fill.solid()
    cross_h.fill.fore_color.rgb = RGBColor(255, 255, 255)
    cross_h.line.fill.background()

    cross_v = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(11.75), Inches(0.05),
        Pt(5), Inches(0.5)
    )
    cross_v.fill.solid()
    cross_v.fill.fore_color.rgb = RGBColor(255, 255, 255)
    cross_v.line.fill.background()

    prs.save('ppt/pptx_templates/medical_clean.pptx')
    print("✓ 医疗简洁模板已创建: medical_clean.pptx")


def create_all_templates():
    """创建所有新模板"""
    os.makedirs('ppt/pptx_templates', exist_ok=True)

    print("\n创建新模板...")
    print("=" * 50)

    # 原有模板
    create_dark_theme_template()
    create_minimalist_template()
    create_chinese_style_template()
    create_gradient_blue_template()
    create_tech_modern_template()
    create_warm_sunset_template()

    # 新增模板
    create_corporate_blue_template()
    create_nature_green_template()
    create_elegant_purple_template()
    create_startup_neon_template()
    create_education_light_template()
    create_medical_clean_template()

    print("=" * 50)
    print("所有新模板创建完成！共 12 个模板\n")


if __name__ == "__main__":
    create_all_templates()
