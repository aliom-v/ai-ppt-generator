"""统一的 PPT 构建器 - 增强版（支持更多页面类型、跨平台字体、统一样式、动画效果）"""
import os
import threading
from typing import Optional, List, Dict, Tuple, NamedTuple
from dataclasses import dataclass
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.shapes.base import BaseShape

from core.ppt_plan import PptPlan, Slide
from utils.logger import get_logger
from ppt.styles import (
    FONT_CN, FONT_EN, ColorTheme,
    CARD_BG_COLORS, ACCENT_COLORS, GRADIENT_COLORS
)


# ==================== 智能内容适配模块 ====================

@dataclass
class ContentMetrics:
    """内容度量指标"""
    total_chars: int = 0           # 总字符数
    max_item_chars: int = 0        # 最长项字符数
    avg_item_chars: float = 0      # 平均项字符数
    item_count: int = 0            # 项目数量
    has_long_items: bool = False   # 是否有长文本项
    density: str = "normal"        # 内容密度: sparse, normal, dense


@dataclass
class LayoutParams:
    """布局参数"""
    font_size: int = 14            # 字体大小
    spacing: float = 1.1           # 行/项间距
    margin: float = 0.4            # 边距
    item_height: float = 1.0       # 每项高度
    max_items: int = 5             # 最大显示项数
    use_compact: bool = False      # 是否使用紧凑布局
    title_font_size: int = 26      # 标题字体大小


def analyze_content(items: List[str]) -> ContentMetrics:
    """分析内容特征，返回度量指标

    Args:
        items: 内容项列表（bullets、时间线项等）

    Returns:
        ContentMetrics 度量指标
    """
    if not items:
        return ContentMetrics()

    filtered = [item for item in items if item and item.strip()]
    if not filtered:
        return ContentMetrics()

    total_chars = sum(len(item) for item in filtered)
    max_item_chars = max(len(item) for item in filtered)
    avg_item_chars = total_chars / len(filtered)
    item_count = len(filtered)

    # 判断是否有长文本项（超过60字符）
    has_long_items = max_item_chars > 60

    # 判断内容密度
    if total_chars < 100 and item_count <= 3:
        density = "sparse"
    elif total_chars > 400 or item_count > 5 or max_item_chars > 100:
        density = "dense"
    else:
        density = "normal"

    return ContentMetrics(
        total_chars=total_chars,
        max_item_chars=max_item_chars,
        avg_item_chars=avg_item_chars,
        item_count=item_count,
        has_long_items=has_long_items,
        density=density
    )


def calculate_layout_params(metrics: ContentMetrics, slide_type: str = "bullets") -> LayoutParams:
    """根据内容度量计算最优布局参数

    Args:
        metrics: 内容度量指标
        slide_type: 页面类型

    Returns:
        LayoutParams 布局参数
    """
    params = LayoutParams()

    # 基于内容密度调整
    if metrics.density == "sparse":
        params.font_size = 16
        params.spacing = 1.3
        params.title_font_size = 28
        params.use_compact = False
    elif metrics.density == "dense":
        params.font_size = 11
        params.spacing = 0.95
        params.title_font_size = 24
        params.use_compact = True
    else:
        params.font_size = 14
        params.spacing = 1.1
        params.title_font_size = 26
        params.use_compact = False

    # 根据最长项调整字体
    if metrics.max_item_chars > 120:
        params.font_size = min(params.font_size, 10)
    elif metrics.max_item_chars > 80:
        params.font_size = min(params.font_size, 11)
    elif metrics.max_item_chars > 50:
        params.font_size = min(params.font_size, 12)

    # 根据项目数量调整
    if metrics.item_count > 6:
        params.max_items = 6
        params.spacing = min(params.spacing, 0.9)
    elif metrics.item_count > 5:
        params.max_items = 5
        params.spacing = min(params.spacing, 1.0)
    else:
        params.max_items = 5

    # 计算每项高度（基于可用空间）
    available_height = 5.5  # 内容区域高度（英寸）
    params.item_height = min(
        params.spacing,
        available_height / max(metrics.item_count, 1)
    )

    # 特定页面类型调整
    if slide_type == "timeline":
        params.font_size = min(params.font_size, 11)
        params.max_items = min(params.max_items, 5)
    elif slide_type == "comparison":
        params.font_size = min(params.font_size, 12)
    elif slide_type == "quote":
        # 引用页需要更大的字体
        if metrics.total_chars < 100:
            params.font_size = 22
        elif metrics.total_chars < 150:
            params.font_size = 19
        else:
            params.font_size = 16

    return params


def smart_text_split(text: str, max_chars: int = 80) -> List[str]:
    """智能文本分割，在合适的位置断行

    Args:
        text: 原始文本
        max_chars: 每行最大字符数

    Returns:
        分割后的文本行列表
    """
    if len(text) <= max_chars:
        return [text]

    lines = []
    current = text

    # 优先在标点符号处断行
    breakpoints = ['。', '；', '，', '、', '：', ' ', '.', ',', ';', ':']

    while len(current) > max_chars:
        # 在 max_chars 附近寻找断点
        best_break = max_chars
        for i in range(max_chars, max(max_chars - 20, 0), -1):
            if i < len(current) and current[i] in breakpoints:
                best_break = i + 1
                break

        lines.append(current[:best_break].strip())
        current = current[best_break:].strip()

    if current:
        lines.append(current)

    return lines


def calculate_font_size_for_text(text: str, max_width: float, base_size: int = 14) -> int:
    """根据文本长度和可用宽度计算合适的字体大小

    Args:
        text: 文本内容
        max_width: 最大可用宽度（英寸）
        base_size: 基础字体大小

    Returns:
        推荐的字体大小
    """
    text_len = len(text)
    # 估算：中文字符约 0.15 英寸宽（14pt），英文约 0.08 英寸
    # 简化处理：假设平均 0.12 英寸/字符
    chars_per_inch = 7  # 14pt 字体大约每英寸 7 个中文字符

    estimated_width = text_len / chars_per_inch
    if estimated_width <= max_width:
        return base_size

    # 需要缩小字体
    scale = max_width / estimated_width
    new_size = int(base_size * scale)
    return max(new_size, 9)  # 最小 9pt

# 尝试导入动画模块
try:
    from ppt.animations import (
        get_animation_builder, apply_animations_to_slide,
        add_simple_transition, AnimationConfig, TransitionConfig
    )
    ANIMATION_AVAILABLE = True
except ImportError:
    ANIMATION_AVAILABLE = False
    logger_temp = get_logger("ppt_builder")
    logger_temp.warning("动画模块不可用，将跳过动画效果")

logger = get_logger("ppt_builder")

# 线程安全的样式计数器（解决多用户并发问题）
_thread_local = threading.local()

# 尝试导入图片搜索模块
try:
    from utils.image_search import search_and_download_image, download_images_parallel
    IMAGE_SEARCH_AVAILABLE = True
except ImportError:
    IMAGE_SEARCH_AVAILABLE = False

# 尝试导入图片增强模块
try:
    from utils.image_enhancer import (
        enhance_image, validate_image, resize_for_ppt,
        EnhanceConfig, get_preset_config, PIL_AVAILABLE as ENHANCE_AVAILABLE
    )
except ImportError:
    ENHANCE_AVAILABLE = False


def _get_bullets_counter() -> int:
    """获取当前线程的 bullets 样式计数器"""
    if not hasattr(_thread_local, 'bullets_counter'):
        _thread_local.bullets_counter = 0
    return _thread_local.bullets_counter


def _increment_bullets_counter() -> int:
    """递增并返回当前线程的 bullets 样式计数器"""
    if not hasattr(_thread_local, 'bullets_counter'):
        _thread_local.bullets_counter = 0
    counter = _thread_local.bullets_counter
    _thread_local.bullets_counter += 1
    return counter


def _reset_bullets_counter():
    """重置当前线程的 bullets 样式计数器"""
    _thread_local.bullets_counter = 0


def build_ppt_from_plan(
    plan: PptPlan,
    template_path: Optional[str],
    output_path: str,
    auto_download_images: bool = False,
    enable_animations: bool = True,
    enhance_images: bool = True,
    image_enhance_preset: str = "default"
) -> None:
    """根据 PptPlan 生成 PPTX 文件

    Args:
        plan: PPT 计划对象
        template_path: 模板路径
        output_path: 输出路径
        auto_download_images: 是否自动下载图片
        enable_animations: 是否启用动画效果（默认 True）
        enhance_images: 是否增强图片（默认 True）
        image_enhance_preset: 图片增强预设（default, vivid, soft, sharp, rounded, shadow）
    """
    _reset_bullets_counter()  # 重置线程安全计数器

    use_template = template_path and os.path.exists(template_path) and template_path.endswith('.pptx')

    if use_template:
        logger.info(f"使用模板: {template_path}")
        prs = Presentation(template_path)
    else:
        logger.info("使用默认样式")
        prs = Presentation()

    # 预下载所有图片（并行）
    if auto_download_images and IMAGE_SEARCH_AVAILABLE:
        _predownload_images(plan.slides, enhance_images, image_enhance_preset)

    # 初始化动画构建器
    anim_builder = None
    if enable_animations and ANIMATION_AVAILABLE:
        anim_builder = get_animation_builder(True)
        logger.info("动画效果已启用")

    # 创建封面页
    cover_slide = _create_title_slide(prs, plan.title, plan.subtitle)
    if anim_builder and cover_slide:
        add_simple_transition(cover_slide, "fade", 1000)

    # 统计动画效果
    animation_count = 0
    transition_count = 1 if anim_builder else 0

    # 创建内容页
    for slide_data in plan.slides:
        slide_type = slide_data.slide_type.lower()
        slide_obj = None

        if slide_type == "bullets":
            slide_obj = _create_bullets_slide(prs, slide_data, anim_builder)
        elif slide_type == "image_with_text":
            slide_obj = _create_image_with_text_slide(prs, slide_data, anim_builder)
        elif slide_type == "two_column":
            slide_obj = _create_two_column_slide(prs, slide_data, anim_builder)
        elif slide_type == "timeline":
            slide_obj = _create_timeline_slide(prs, slide_data, anim_builder)
        elif slide_type == "comparison":
            slide_obj = _create_comparison_slide(prs, slide_data, anim_builder)
        elif slide_type == "quote":
            slide_obj = _create_quote_slide(prs, slide_data, anim_builder)
        elif slide_type == "ending":
            slide_obj = _create_ending_slide(prs, slide_data, anim_builder)
        else:
            slide_obj = _create_bullets_slide(prs, slide_data, anim_builder)

        # 统计动画
        if anim_builder and slide_obj:
            transition_count += 1

    # 确保输出目录存在
    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)

    prs.save(output_path)

    if enable_animations and ANIMATION_AVAILABLE:
        logger.info(f"PPT 已保存: {output_path}（含 {transition_count} 个页面切换效果）")
    else:
        logger.info(f"PPT 已保存: {output_path}")


def _predownload_images(
    slides: List[Slide],
    enhance: bool = True,
    preset: str = "default"
) -> None:
    """预下载所有需要的图片（并行）并进行增强处理

    Args:
        slides: 幻灯片列表
        enhance: 是否增强图片
        preset: 增强预设名称
    """
    keywords = []
    keyword_to_slides: dict = {}

    for slide in slides:
        if slide.slide_type.lower() == "image_with_text":
            if slide.image_keyword and not (slide.image_path and os.path.exists(slide.image_path)):
                kw = slide.image_keyword
                if kw not in keyword_to_slides:
                    keywords.append(kw)
                    keyword_to_slides[kw] = []
                keyword_to_slides[kw].append(slide)

    if not keywords:
        return

    logger.info(f"并行下载 {len(keywords)} 张图片...")
    results = download_images_parallel(keywords)

    # 处理下载结果
    enhanced_count = 0
    for keyword, path in results.items():
        if path and keyword in keyword_to_slides:
            # 图片增强处理
            final_path = path
            if enhance and ENHANCE_AVAILABLE:
                try:
                    # 获取预设配置
                    config = get_preset_config(preset)

                    # 先验证图片质量
                    report = validate_image(path)
                    if report.is_valid:
                        # 调整大小（如果需要）
                        resized_path = resize_for_ppt(path)
                        if resized_path:
                            path = resized_path

                        # 应用增强效果
                        enhanced_path = enhance_image(path, config=config)
                        if enhanced_path:
                            final_path = enhanced_path
                            enhanced_count += 1
                    else:
                        logger.warning(f"图片质量问题: {keyword} - {report.issues}")
                except Exception as e:
                    logger.warning(f"图片增强失败: {keyword} - {e}")

            # 更新所有使用该关键词的幻灯片
            for slide in keyword_to_slides[keyword]:
                slide.image_path = final_path
            logger.debug(f"下载完成: {keyword}")

    if enhanced_count > 0:
        logger.info(f"已增强 {enhanced_count} 张图片")


def _set_font(text_frame, font_name: str = None, font_size: int = None, bold: bool = False, color: RGBColor = None):
    """设置文本框字体（跨平台兼容）"""
    if font_name is None:
        font_name = FONT_CN
    
    if text_frame and text_frame.paragraphs:
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = font_name
            if font_size:
                paragraph.font.size = Pt(font_size)
            if bold:
                paragraph.font.bold = True
            if color:
                paragraph.font.color.rgb = color
            for run in paragraph.runs:
                run.font.name = font_name
                if font_size:
                    run.font.size = Pt(font_size)
                if bold:
                    run.font.bold = True
                if color:
                    run.font.color.rgb = color


def _add_header_decoration(slide, prs: Presentation):
    """添加页面顶部装饰条"""
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Pt(8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = ColorTheme.PRIMARY
    header.line.fill.background()


def _add_footer_decoration(slide, prs: Presentation):
    """添加页面底部装饰"""
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), prs.slide_height - Pt(4),
        prs.slide_width, Pt(4)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = ColorTheme.SECONDARY
    footer.line.fill.background()


def _add_side_accent(slide, prs: Presentation):
    """添加左侧装饰条"""
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.8),
        Pt(6), Inches(0.6)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ColorTheme.ACCENT
    accent.line.fill.background()


def _create_title_slide(prs: Presentation, title: str, subtitle: str):
    """创建封面页

    Returns:
        创建的幻灯片对象
    """
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    try:
        slide.shapes.title.text = title
        _set_font(slide.shapes.title.text_frame, font_size=44, bold=True)
    except (AttributeError, IndexError):
        pass  # 模板可能没有标题占位符

    try:
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
            _set_font(slide.placeholders[1].text_frame, font_size=24)
    except (AttributeError, IndexError):
        pass  # 模板可能没有副标题占位符

    return slide


def _get_slide_width_inches(prs: Presentation) -> float:
    """获取幻灯片宽度（英寸）"""
    return prs.slide_width.inches if hasattr(prs.slide_width, 'inches') else prs.slide_width / 914400


def _create_bullets_slide(prs: Presentation, slide_data: Slide, anim_builder=None):
    """创建要点页 - 多种样式交替

    Returns:
        创建的幻灯片对象
    """
    bullets = [b for b in (slide_data.bullets or []) if b]
    num_bullets = len(bullets)

    if num_bullets == 0:
        return None

    # 使用线程安全的计数器选择样式（4种样式交替）
    style = _increment_bullets_counter() % 4

    if style == 0:
        slide = _create_bullets_style_cards(prs, slide_data, bullets)
    elif style == 1:
        slide = _create_bullets_style_list(prs, slide_data, bullets)
    elif style == 2:
        slide = _create_bullets_style_icons(prs, slide_data, bullets)
    else:
        slide = _create_bullets_style_gradient(prs, slide_data, bullets)

    # 添加页面切换效果
    if anim_builder and slide and ANIMATION_AVAILABLE:
        add_simple_transition(slide, "push", 500)

    return slide


def _create_bullets_style_cards(prs: Presentation, slide_data: Slide, bullets: List[str]):
    """样式1: 卡片式布局 - 居中（智能内容适配）

    Returns:
        创建的幻灯片对象
    """
    layout_idx = min(6, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    slide_w = _get_slide_width_inches(prs)
    _add_header_decoration(slide, prs)

    # 智能内容分析
    metrics = analyze_content(bullets)
    layout = calculate_layout_params(metrics, "bullets")

    # 标题 - 居中，根据内容密度调整
    margin = 0.4
    title_box = slide.shapes.add_textbox(Inches(margin), Inches(0.25), Inches(slide_w - 2*margin), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(title_frame, font_size=layout.title_font_size, bold=True, color=ColorTheme.TEXT_DARK)

    # 使用智能计算的布局参数
    content_width = slide_w - 1.0  # 左右各留0.5英寸
    start_x = 0.5
    spacing = layout.item_height
    max_items = layout.max_items

    for i, bullet in enumerate(bullets[:max_items]):
        y = 1.0 + i * spacing
        _draw_bullet_card_horizontal(
            slide, start_x, y, content_width, spacing - 0.1, i + 1, bullet,
            CARD_BG_COLORS[i % len(CARD_BG_COLORS)], ACCENT_COLORS[i % len(ACCENT_COLORS)],
            font_size=layout.font_size
        )

    return slide


def _create_bullets_style_list(prs: Presentation, slide_data: Slide, bullets: List[str]):
    """样式2: 简洁列表式 - 居中（智能内容适配）

    Returns:
        创建的幻灯片对象
    """
    layout_idx = min(6, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    slide_w = _get_slide_width_inches(prs)
    _add_header_decoration(slide, prs)

    # 智能内容分析
    metrics = analyze_content(bullets)
    layout = calculate_layout_params(metrics, "bullets")

    # 标题 - 居中
    margin = 0.4
    title_box = slide.shapes.add_textbox(Inches(margin), Inches(0.25), Inches(slide_w - 2*margin), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(title_frame, font_size=layout.title_font_size, bold=True, color=ColorTheme.TEXT_DARK)

    # 居中计算
    content_width = slide_w - 1.0
    start_x = 0.5

    # 内容区域背景 - 居中
    content_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(start_x), Inches(0.95), Inches(content_width), Inches(5.8)
    )
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
    content_bg.line.fill.background()

    # 使用智能计算的布局参数
    font_size = layout.font_size
    spacing = layout.item_height
    max_items = layout.max_items
    colors = ACCENT_COLORS[:5]

    for i, bullet in enumerate(bullets[:max_items]):
        y = 1.1 + i * spacing

        # 序号圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(start_x + 0.15), Inches(y + 0.08), Inches(0.32), Inches(0.32))
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors[i % len(colors)]
        dot.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(start_x + 0.15), Inches(y + 0.1), Inches(0.32), Inches(0.32))
        num_frame = num_box.text_frame
        num_frame.text = str(i + 1)
        num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _set_font(num_frame, font_size=12, bold=True, color=ColorTheme.WHITE)

        # 文字
        text_box = slide.shapes.add_textbox(Inches(start_x + 0.6), Inches(y), Inches(content_width - 0.75), Inches(spacing))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = bullet
        _set_font(text_frame, font_size=font_size, color=ColorTheme.TEXT_DARK)

    return slide


def _create_bullets_style_icons(prs: Presentation, slide_data: Slide, bullets: List[str]):
    """样式3: 图标式布局（带大序号）- 智能内容适配

    Returns:
        创建的幻灯片对象
    """
    layout_idx = min(6, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    slide_w = _get_slide_width_inches(prs)
    _add_header_decoration(slide, prs)

    # 智能内容分析
    metrics = analyze_content(bullets)
    layout = calculate_layout_params(metrics, "bullets")

    # 标题 - 居中
    margin = 0.4
    title_box = slide.shapes.add_textbox(Inches(margin), Inches(0.25), Inches(slide_w - 2*margin), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(title_frame, font_size=layout.title_font_size, bold=True, color=ColorTheme.TEXT_DARK)

    colors = [ColorTheme.PRIMARY, ColorTheme.SUCCESS, ColorTheme.ACCENT,
              RGBColor(156, 39, 176), RGBColor(0, 188, 212)]

    # 使用智能计算的布局参数
    content_width = slide_w - 1.0
    start_x = 0.5
    spacing = layout.item_height
    max_items = layout.max_items
    font_size = layout.font_size

    for i, bullet in enumerate(bullets[:max_items]):
        y = 1.0 + i * spacing

        # 序号方块
        num_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(start_x), Inches(y), Inches(0.65), Inches(spacing - 0.15)
        )
        num_bg.fill.solid()
        num_bg.fill.fore_color.rgb = colors[i % len(colors)]
        num_bg.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(start_x), Inches(y + 0.12), Inches(0.65), Inches(0.5))
        num_frame = num_box.text_frame
        num_frame.text = str(i + 1)
        num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _set_font(num_frame, font_size=22, bold=True, color=ColorTheme.WHITE)

        # 内容背景
        content_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(start_x + 0.75), Inches(y), Inches(content_width - 0.75), Inches(spacing - 0.15)
        )
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
        content_bg.line.fill.background()

        # 内容文字 - 使用智能字体大小
        text_box = slide.shapes.add_textbox(Inches(start_x + 0.85), Inches(y + 0.08), Inches(content_width - 0.95), Inches(spacing - 0.2))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = bullet
        _set_font(text_frame, font_size=font_size, color=ColorTheme.TEXT_DARK)

    return slide


def _create_bullets_style_gradient(prs: Presentation, slide_data: Slide, bullets: List[str]):
    """样式4: 渐变色条式 - 智能内容适配

    Returns:
        创建的幻灯片对象
    """
    layout_idx = min(6, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    slide_w = _get_slide_width_inches(prs)
    _add_header_decoration(slide, prs)

    # 智能内容分析
    metrics = analyze_content(bullets)
    layout = calculate_layout_params(metrics, "bullets")

    # 标题 - 居中
    margin = 0.4
    title_box = slide.shapes.add_textbox(Inches(margin), Inches(0.25), Inches(slide_w - 2*margin), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(title_frame, font_size=layout.title_font_size, bold=True, color=ColorTheme.TEXT_DARK)

    # 渐变色系
    gradient_colors = GRADIENT_COLORS
    bg_colors = CARD_BG_COLORS[:5]

    # 使用智能计算的布局参数
    content_width = slide_w - 1.0
    start_x = 0.5
    spacing = layout.item_height
    max_items = layout.max_items
    font_size = layout.font_size

    for i, bullet in enumerate(bullets[:max_items]):
        y = 1.0 + i * spacing

        # 背景条 - 居中
        bar_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(start_x), Inches(y), Inches(content_width), Inches(spacing - 0.1)
        )
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = bg_colors[i % len(bg_colors)]
        bar_bg.line.fill.background()

        # 左侧色条
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(start_x), Inches(y), Pt(8), Inches(spacing - 0.1)
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = gradient_colors[i % len(gradient_colors)]
        left_bar.line.fill.background()

        # 序号
        num_box = slide.shapes.add_textbox(Inches(start_x + 0.2), Inches(y + 0.05), Inches(0.4), Inches(spacing - 0.2))
        num_frame = num_box.text_frame
        num_frame.text = str(i + 1)
        _set_font(num_frame, font_size=18, bold=True, color=gradient_colors[i % len(gradient_colors)])

        # 内容 - 使用智能字体大小
        text_box = slide.shapes.add_textbox(Inches(start_x + 0.65), Inches(y + 0.1), Inches(content_width - 0.85), Inches(spacing - 0.2))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = bullet
        _set_font(text_frame, font_size=font_size, color=ColorTheme.TEXT_DARK)

    return slide


def _draw_bullet_card(slide, x, y, width, height, num, text, bg_color, accent_color, horizontal=False):
    """绘制单个要点卡片（垂直布局）"""
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.fill.background()
    
    # 序号圆圈
    circle_size = 0.4
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + 0.15), Inches(y + 0.15),
        Inches(circle_size), Inches(circle_size)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()
    
    # 序号文字
    num_box = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.18),
        Inches(circle_size), Inches(circle_size)
    )
    num_frame = num_box.text_frame
    num_frame.text = str(num)
    num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(num_frame, font_size=14, bold=True, color=ColorTheme.WHITE)
    
    # 内容文字
    text_len = len(text)
    if horizontal:
        # 横向布局：序号在左，文字在右
        if text_len > 80:
            font_size = 11
        elif text_len > 50:
            font_size = 12
        else:
            font_size = 13
        text_box = slide.shapes.add_textbox(
            Inches(x + 0.7), Inches(y + 0.2),
            Inches(width - 0.9), Inches(height - 0.4)
        )
    else:
        # 垂直布局：序号在上，文字在下
        if text_len > 100:
            font_size = 11
        elif text_len > 60:
            font_size = 12
        else:
            font_size = 13
        text_box = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + 0.7),
            Inches(width - 0.3), Inches(height - 0.9)
        )
    
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = text
    _set_font(text_frame, font_size=font_size, color=ColorTheme.TEXT_DARK)


def _draw_bullet_card_horizontal(slide, x, y, width, height, num, text, bg_color, accent_color, font_size: int = None):
    """绘制横向要点卡片（列表式）

    Args:
        font_size: 字体大小，如果为 None 则根据文本长度自动计算
    """
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.fill.background()

    # 左侧色条
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y),
        Pt(6), Inches(height)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent_color
    accent_bar.line.fill.background()

    # 序号圆圈
    circle_size = 0.35
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + 0.15), Inches(y + (height - 0.35) / 2),
        Inches(circle_size), Inches(circle_size)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()

    # 序号文字
    num_box = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + (height - 0.35) / 2 + 0.02),
        Inches(circle_size), Inches(circle_size)
    )
    num_frame = num_box.text_frame
    num_frame.text = str(num)
    num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(num_frame, font_size=12, bold=True, color=ColorTheme.WHITE)

    # 内容文字 - 如果没有指定 font_size，则根据文本长度计算
    if font_size is None:
        text_len = len(text)
        if text_len > 80:
            font_size = 12
        elif text_len > 50:
            font_size = 13
        else:
            font_size = 14

    text_box = slide.shapes.add_textbox(
        Inches(x + 0.6), Inches(y + 0.1),
        Inches(width - 0.8), Inches(height - 0.2)
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = text
    text_frame.paragraphs[0].vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_font(text_frame, font_size=font_size, color=ColorTheme.TEXT_DARK)


def _create_image_with_text_slide(prs: Presentation, slide_data: Slide, anim_builder=None):
    """创建图文混排页 - 智能内容适配

    Returns:
        创建的幻灯片对象
    """
    layout_idx = min(6, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    slide_w = _get_slide_width_inches(prs)
    _add_header_decoration(slide, prs)

    # 智能内容分析
    text_content = slide_data.text or "（图片说明）"
    metrics = analyze_content([text_content])
    layout = calculate_layout_params(metrics, "image_with_text")

    # 标题 - 居中
    margin = 0.4
    title_box = slide.shapes.add_textbox(Inches(margin), Inches(0.25), Inches(slide_w - 2*margin), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(title_frame, font_size=layout.title_font_size, bold=True, color=ColorTheme.TEXT_DARK)

    # 计算居中位置
    content_width = slide_w - 1.0
    start_x = 0.5
    img_width = content_width * 0.45
    text_width = content_width * 0.52
    gap = content_width * 0.03

    # 图片或占位符 - 左侧
    img_left = start_x
    if slide_data.image_path and os.path.exists(slide_data.image_path):
        try:
            slide.shapes.add_picture(
                slide_data.image_path,
                Inches(img_left), Inches(1.1),
                width=Inches(img_width)
            )
        except Exception as e:
            logger.warning(f"无法插入图片: {e}")
            _add_image_placeholder(slide, slide_data, img_left, img_width)
    else:
        _add_image_placeholder(slide, slide_data, img_left, img_width)

    # 文字说明（带背景框）- 右侧
    text_left = img_left + img_width + gap
    text_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(text_left), Inches(1.1),
        Inches(text_width), Inches(5.5)
    )
    text_bg.fill.solid()
    text_bg.fill.fore_color.rgb = ColorTheme.BG_LIGHT
    text_bg.line.fill.background()

    # 使用智能计算的字体大小
    font_size = layout.font_size

    text_box = slide.shapes.add_textbox(Inches(text_left + 0.15), Inches(1.25), Inches(text_width - 0.3), Inches(5.2))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = text_content
    _set_font(text_frame, font_size=font_size, color=ColorTheme.TEXT_DARK)

    # 添加页面切换效果
    if anim_builder and ANIMATION_AVAILABLE:
        add_simple_transition(slide, "fade", 700)

    return slide


def _add_image_placeholder(slide, slide_data: Slide, left: float = 0.75, width: float = 4.0) -> None:
    """添加图片占位符"""
    # 占位符背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(1.2),
        Inches(width), Inches(5.3)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = ColorTheme.BG_LIGHT
    bg.line.color.rgb = RGBColor(200, 200, 200)
    bg.line.width = Pt(1)
    
    box = slide.shapes.add_textbox(Inches(left), Inches(1.2), Inches(width), Inches(5.3))
    frame = box.text_frame
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    text = f"🖼️ 图片占位符\n\n关键词: {slide_data.image_keyword}" if slide_data.image_keyword else "🖼️ 图片占位符\n\n请在此处插入图片"
    
    frame.text = text
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.name = FONT_CN
    p.font.color.rgb = ColorTheme.TEXT_LIGHT


def _create_two_column_slide(prs: Presentation, slide_data: Slide, anim_builder=None):
    """创建双栏布局页 - 智能内容适配

    Returns:
        创建的幻灯片对象
    """
    layout_idx = min(6, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    slide_w = _get_slide_width_inches(prs)
    _add_header_decoration(slide, prs)

    # 分割要点为左右两栏
    bullets = slide_data.bullets or []
    mid = (len(bullets) + 1) // 2
    left_bullets = bullets[:mid]
    right_bullets = bullets[mid:]

    # 智能内容分析
    metrics = analyze_content(bullets)
    layout = calculate_layout_params(metrics, "two_column")

    # 标题 - 居中
    margin = 0.4
    title_box = slide.shapes.add_textbox(Inches(margin), Inches(0.25), Inches(slide_w - 2*margin), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(title_frame, font_size=layout.title_font_size, bold=True, color=ColorTheme.TEXT_DARK)

    # 居中计算
    content_width = slide_w - 1.0
    start_x = 0.5
    col_width = (content_width - 0.3) / 2
    gap = 0.3
    font_size = layout.font_size

    # 左栏背景
    left_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(start_x), Inches(1.0),
        Inches(col_width), Inches(5.6)
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = ColorTheme.BG_LIGHT
    left_bg.line.fill.background()

    # 左栏内容
    left_box = slide.shapes.add_textbox(Inches(start_x + 0.12), Inches(1.15), Inches(col_width - 0.24), Inches(5.3))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        p = left_frame.paragraphs[0] if i == 0 else left_frame.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = FONT_CN
        p.font.size = Pt(font_size)
        p.font.color.rgb = ColorTheme.TEXT_DARK
        p.space_after = Pt(10)

    # 右栏背景
    right_x = start_x + col_width + gap
    right_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(right_x), Inches(1.0),
        Inches(col_width), Inches(5.6)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(240, 248, 255)  # 浅蓝色
    right_bg.line.fill.background()

    # 右栏内容
    right_box = slide.shapes.add_textbox(Inches(right_x + 0.12), Inches(1.15), Inches(col_width - 0.24), Inches(5.3))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        p = right_frame.paragraphs[0] if i == 0 else right_frame.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = FONT_CN
        p.font.size = Pt(font_size)
        p.font.color.rgb = ColorTheme.TEXT_DARK
        p.space_after = Pt(10)

    # 添加页面切换效果
    if anim_builder and ANIMATION_AVAILABLE:
        add_simple_transition(slide, "split", 600)

    return slide


def _create_timeline_slide(prs: Presentation, slide_data: Slide, anim_builder=None):
    """创建时间线页 - 智能内容适配

    Returns:
        创建的幻灯片对象
    """
    layout_idx = min(6, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    slide_w = _get_slide_width_inches(prs)
    _add_header_decoration(slide, prs)

    bullets = slide_data.bullets or []

    # 智能内容分析
    metrics = analyze_content(bullets)
    layout = calculate_layout_params(metrics, "timeline")

    # 标题 - 居中
    margin = 0.4
    title_box = slide.shapes.add_textbox(Inches(margin), Inches(0.25), Inches(slide_w - 2*margin), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(title_frame, font_size=layout.title_font_size, bold=True, color=ColorTheme.TEXT_DARK)

    num_items = min(len(bullets), layout.max_items)

    if num_items == 0:
        return slide

    # 绘制时间线主轴 - 居中
    content_width = slide_w - 1.0
    start_x = 0.5
    line_y = Inches(3.75)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(start_x), line_y,
        Inches(content_width), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ColorTheme.PRIMARY
    line.line.fill.background()

    # 根据数量计算间距 - 基于实际宽度
    item_spacing = content_width / (num_items + 1)
    positions = [start_x + item_spacing * (i + 1) for i in range(num_items)]

    colors = [ColorTheme.PRIMARY, ColorTheme.SECONDARY, ColorTheme.ACCENT, ColorTheme.SUCCESS, ColorTheme.WARNING]

    # 根据内容度量决定卡片宽度
    if metrics.max_item_chars > 50:
        card_width = 1.8
    elif metrics.max_item_chars > 35:
        card_width = 1.6
    else:
        card_width = 1.5

    for i, bullet in enumerate(bullets[:num_items]):
        x = Inches(positions[i])

        # 圆点
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x - Inches(0.12), line_y - Inches(0.1),
            Inches(0.24), Inches(0.24)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors[i % len(colors)]
        circle.line.color.rgb = ColorTheme.WHITE
        circle.line.width = Pt(2)

        # 序号
        num_box = slide.shapes.add_textbox(x - Inches(0.08), line_y - Inches(0.06), Inches(0.16), Inches(0.16))
        num_frame = num_box.text_frame
        num_frame.text = str(i + 1)
        num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _set_font(num_frame, font_size=10, bold=True, color=ColorTheme.WHITE)

        # 文字卡片（交替上下）
        is_top = i % 2 == 0
        card_y = Inches(1.0) if is_top else Inches(4.1)
        card_height = 2.4

        # 卡片背景
        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x - Inches(card_width / 2), card_y,
            Inches(card_width), Inches(card_height)
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = colors[i % len(colors)]
        card_bg.line.fill.background()

        # 连接线
        conn_y = Inches(3.4) if is_top else line_y + Pt(4)
        conn_height = Inches(0.35) if is_top else Inches(0.35)
        conn = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x - Pt(2), conn_y,
            Pt(4), conn_height
        )
        conn.fill.solid()
        conn.fill.fore_color.rgb = colors[i % len(colors)]
        conn.line.fill.background()

        # 卡片文字 - 使用智能字体大小
        font_size = layout.font_size
        # 时间线卡片特殊处理：进一步根据单项长度微调
        text_len = len(bullet)
        if text_len > 60:
            font_size = min(font_size, 8)
        elif text_len > 45:
            font_size = min(font_size, 9)
        elif text_len > 30:
            font_size = min(font_size, 10)

        text_box = slide.shapes.add_textbox(
            x - Inches(card_width / 2 - 0.08), card_y + Inches(0.08),
            Inches(card_width - 0.16), Inches(card_height - 0.16)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = bullet
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _set_font(text_frame, font_size=font_size, color=ColorTheme.WHITE)

    # 添加页面切换效果
    if anim_builder and ANIMATION_AVAILABLE:
        add_simple_transition(slide, "wipe", 800)

    return slide


def _create_comparison_slide(prs: Presentation, slide_data: Slide, anim_builder=None):
    """创建对比页（左右对比）- 智能内容适配

    Returns:
        创建的幻灯片对象
    """
    layout_idx = min(6, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    slide_w = _get_slide_width_inches(prs)
    _add_header_decoration(slide, prs)

    bullets = slide_data.bullets or []
    mid = (len(bullets) + 1) // 2
    left_items = bullets[:mid]
    right_items = bullets[mid:]

    # 智能内容分析
    metrics = analyze_content(bullets)
    layout = calculate_layout_params(metrics, "comparison")

    # 标题 - 居中
    margin = 0.4
    title_box = slide.shapes.add_textbox(Inches(margin), Inches(0.25), Inches(slide_w - 2*margin), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(title_frame, font_size=layout.title_font_size, bold=True, color=ColorTheme.TEXT_DARK)

    # 居中计算
    content_width = slide_w - 1.0
    start_x = 0.5
    col_width = (content_width - 0.3) / 2
    gap = 0.3
    right_x = start_x + col_width + gap
    font_size = layout.font_size

    # 左侧标题框
    left_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(start_x), Inches(0.9),
        Inches(col_width), Inches(0.45)
    )
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = ColorTheme.PRIMARY
    left_header.line.fill.background()
    left_header.text_frame.text = "方案 A"
    left_header.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(left_header.text_frame, font_size=15, bold=True, color=ColorTheme.WHITE)

    # 左侧内容背景
    left_content_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(start_x), Inches(1.4),
        Inches(col_width), Inches(5.2)
    )
    left_content_bg.fill.solid()
    left_content_bg.fill.fore_color.rgb = RGBColor(240, 248, 255)
    left_content_bg.line.fill.background()

    # 左侧内容
    left_box = slide.shapes.add_textbox(Inches(start_x + 0.15), Inches(1.65), Inches(col_width - 0.3), Inches(4.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        p = left_frame.paragraphs[0] if i == 0 else left_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT_CN
        p.font.size = Pt(font_size)
        p.font.color.rgb = ColorTheme.TEXT_DARK
        p.space_after = Pt(10)

    # 右侧标题框
    right_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(right_x), Inches(0.9),
        Inches(col_width), Inches(0.45)
    )
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = ColorTheme.ACCENT
    right_header.line.fill.background()
    right_header.text_frame.text = "方案 B"
    right_header.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(right_header.text_frame, font_size=15, bold=True, color=ColorTheme.WHITE)

    # 右侧内容背景
    right_content_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(right_x), Inches(1.4),
        Inches(col_width), Inches(5.2)
    )
    right_content_bg.fill.solid()
    right_content_bg.fill.fore_color.rgb = RGBColor(255, 248, 240)
    right_content_bg.line.fill.background()

    # 右侧内容
    right_box = slide.shapes.add_textbox(Inches(right_x + 0.12), Inches(1.55), Inches(col_width - 0.24), Inches(4.9))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        p = right_frame.paragraphs[0] if i == 0 else right_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT_CN
        p.font.size = Pt(font_size)
        p.font.color.rgb = ColorTheme.TEXT_DARK
        p.space_after = Pt(10)

    # 添加页面切换效果
    if anim_builder and ANIMATION_AVAILABLE:
        add_simple_transition(slide, "split", 600)

    return slide


def _create_quote_slide(prs: Presentation, slide_data: Slide, anim_builder=None):
    """创建引用/名言页 - 智能内容适配

    Returns:
        创建的幻灯片对象
    """
    layout_idx = min(6, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    slide_w = _get_slide_width_inches(prs)
    _add_header_decoration(slide, prs)

    # 智能内容分析
    quote_text = slide_data.text or slide_data.title
    metrics = analyze_content([quote_text])
    layout = calculate_layout_params(metrics, "quote")

    # 居中计算
    content_width = slide_w - 1.0
    start_x = 0.5

    # 背景装饰 - 居中
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(start_x), Inches(1.0),
        Inches(content_width), Inches(5.6)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = ColorTheme.BG_LIGHT
    bg_shape.line.fill.background()

    # 大引号装饰 - 居中
    quote_mark = slide.shapes.add_textbox(Inches(start_x + 0.25), Inches(0.8), Inches(1.2), Inches(1.2))
    quote_frame = quote_mark.text_frame
    quote_frame.text = "\u201C"  # 左双引号装饰
    quote_frame.paragraphs[0].font.size = Pt(80)
    quote_frame.paragraphs[0].font.color.rgb = ColorTheme.PRIMARY
    quote_frame.paragraphs[0].font.name = "Georgia"

    # 引用内容 - 使用智能字体大小
    font_size = layout.font_size

    quote_box = slide.shapes.add_textbox(Inches(start_x + 0.4), Inches(2.0), Inches(content_width - 0.8), Inches(3.2))
    quote_text_frame = quote_box.text_frame
    quote_text_frame.word_wrap = True
    quote_text_frame.text = quote_text
    quote_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(quote_text_frame, font_size=font_size, color=ColorTheme.TEXT_DARK)

    # 作者/来源
    if slide_data.subtitle:
        author_box = slide.shapes.add_textbox(Inches(start_x + 0.4), Inches(5.6), Inches(content_width - 0.8), Inches(0.5))
        author_frame = author_box.text_frame
        author_frame.text = f"— {slide_data.subtitle}"
        author_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        _set_font(author_frame, font_size=15, color=ColorTheme.PRIMARY)

    # 添加页面切换效果
    if anim_builder and ANIMATION_AVAILABLE:
        add_simple_transition(slide, "fade", 800)

    return slide


def _create_ending_slide(prs: Presentation, slide_data: Slide, anim_builder=None):
    """创建结束页

    Returns:
        创建的幻灯片对象
    """
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    try:
        slide.shapes.title.text = slide_data.title
        _set_font(slide.shapes.title.text_frame, font_size=44, bold=True)
    except (AttributeError, IndexError):
        pass  # 模板可能没有标题占位符

    try:
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_data.subtitle or ""
            _set_font(slide.placeholders[1].text_frame, font_size=24)
    except (AttributeError, IndexError):
        pass  # 模板可能没有副标题占位符

    # 添加页面切换效果
    if anim_builder and ANIMATION_AVAILABLE:
        add_simple_transition(slide, "fade", 1000)

    return slide
