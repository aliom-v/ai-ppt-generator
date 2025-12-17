"""PPT 动画效果模块 - 基于 OOXML 实现入场动画和页面切换"""
from typing import Optional, List, Literal, Dict, Any
from dataclasses import dataclass, field
from lxml import etree
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.util import Emu

from utils.logger import get_logger

logger = get_logger("animations")

# 命名空间定义
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

# 动画类型
AnimationType = Literal["appear", "fade", "fly_in", "wipe", "zoom", "bounce", "float"]

# 页面切换类型
TransitionType = Literal["fade", "push", "wipe", "split", "reveal", "random"]

# 动画方向
AnimationDirection = Literal["left", "right", "top", "bottom", "top_left", "top_right", "bottom_left", "bottom_right"]


@dataclass
class AnimationConfig:
    """动画配置"""
    anim_type: AnimationType = "appear"
    duration: int = 500  # 毫秒
    delay: int = 0       # 延迟（毫秒）
    direction: AnimationDirection = "left"

    def to_dict(self) -> Dict[str, Any]:
        return {
            "type": self.anim_type,
            "duration": self.duration,
            "delay": self.delay,
            "direction": self.direction
        }


@dataclass
class TransitionConfig:
    """页面切换配置"""
    trans_type: TransitionType = "fade"
    duration: int = 700  # 毫秒

    def to_dict(self) -> Dict[str, Any]:
        return {
            "type": self.trans_type,
            "duration": self.duration
        }


# 页面类型默认动画配置
DEFAULT_SLIDE_ANIMATIONS: Dict[str, Dict[str, Any]] = {
    "bullets": {
        "title": AnimationConfig("fade", 400, 0),
        "items": AnimationConfig("fly_in", 300, 200),
        "item_interval": 150,  # 每项之间的间隔
    },
    "image_with_text": {
        "title": AnimationConfig("fade", 400, 0),
        "image": AnimationConfig("zoom", 500, 200),
        "text": AnimationConfig("fade", 400, 500),
    },
    "two_column": {
        "title": AnimationConfig("fade", 400, 0),
        "left": AnimationConfig("fly_in", 400, 200, "left"),
        "right": AnimationConfig("fly_in", 400, 400, "right"),
    },
    "timeline": {
        "title": AnimationConfig("fade", 400, 0),
        "line": AnimationConfig("wipe", 600, 200),
        "items": AnimationConfig("zoom", 300, 400),
        "item_interval": 200,
    },
    "comparison": {
        "title": AnimationConfig("fade", 400, 0),
        "left": AnimationConfig("fly_in", 400, 200, "left"),
        "right": AnimationConfig("fly_in", 400, 400, "right"),
    },
    "quote": {
        "title": AnimationConfig("fade", 400, 0),
        "quote_mark": AnimationConfig("zoom", 300, 100),
        "text": AnimationConfig("fade", 600, 300),
        "author": AnimationConfig("appear", 300, 800),
    },
    "ending": {
        "title": AnimationConfig("zoom", 600, 0),
        "subtitle": AnimationConfig("fade", 400, 400),
    },
}

# 默认页面切换效果
DEFAULT_TRANSITIONS: Dict[str, TransitionConfig] = {
    "cover": TransitionConfig("fade", 1000),
    "bullets": TransitionConfig("push", 500),
    "image_with_text": TransitionConfig("fade", 700),
    "two_column": TransitionConfig("split", 600),
    "timeline": TransitionConfig("wipe", 800),
    "comparison": TransitionConfig("split", 600),
    "quote": TransitionConfig("fade", 800),
    "ending": TransitionConfig("fade", 1000),
}


def _make_p_tag(tag: str) -> str:
    """创建带命名空间的标签"""
    return f"{{{NSMAP['p']}}}{tag}"


def _make_a_tag(tag: str) -> str:
    """创建带命名空间的标签"""
    return f"{{{NSMAP['a']}}}{tag}"


class AnimationBuilder:
    """PPT 动画构建器

    通过操作 OOXML 为 PPT 添加动画效果。
    支持入场动画、页面切换、顺序动画等。
    """

    def __init__(self, enable_animations: bool = True):
        """初始化动画构建器

        Args:
            enable_animations: 是否启用动画（默认 True）
        """
        self.enabled = enable_animations
        self._shape_counter = 0

    def add_entrance_animation(
        self,
        slide: Slide,
        shape: BaseShape,
        anim_config: Optional[AnimationConfig] = None
    ) -> bool:
        """为形状添加入场动画

        Args:
            slide: 幻灯片对象
            shape: 要添加动画的形状
            anim_config: 动画配置

        Returns:
            是否添加成功
        """
        if not self.enabled:
            return False

        if anim_config is None:
            anim_config = AnimationConfig()

        try:
            # 获取或创建 timing 元素
            timing = self._get_or_create_timing(slide)

            # 获取或创建 tnLst (时间节点列表)
            tn_lst = self._get_or_create_tn_lst(timing)

            # 创建动画节点
            self._create_animation_node(
                tn_lst, shape, anim_config
            )

            self._shape_counter += 1
            return True

        except Exception as e:
            logger.warning(f"添加入场动画失败: {e}")
            return False

    def add_sequence_animation(
        self,
        slide: Slide,
        shapes: List[BaseShape],
        anim_config: Optional[AnimationConfig] = None,
        interval: int = 200
    ) -> int:
        """为多个形状添加顺序动画（逐项出现）

        Args:
            slide: 幻灯片对象
            shapes: 形状列表
            anim_config: 动画配置
            interval: 每个形状之间的延迟（毫秒）

        Returns:
            成功添加动画的形状数量
        """
        if not self.enabled or not shapes:
            return 0

        if anim_config is None:
            anim_config = AnimationConfig("fly_in", 300, 0)

        count = 0
        base_delay = anim_config.delay

        for i, shape in enumerate(shapes):
            config = AnimationConfig(
                anim_type=anim_config.anim_type,
                duration=anim_config.duration,
                delay=base_delay + i * interval,
                direction=anim_config.direction
            )
            if self.add_entrance_animation(slide, shape, config):
                count += 1

        return count

    def add_slide_transition(
        self,
        slide: Slide,
        trans_config: Optional[TransitionConfig] = None
    ) -> bool:
        """添加页面切换效果

        Args:
            slide: 幻灯片对象
            trans_config: 切换配置

        Returns:
            是否添加成功
        """
        if not self.enabled:
            return False

        if trans_config is None:
            trans_config = TransitionConfig()

        try:
            # 获取 slide 的 XML 元素
            sld = slide._element

            # 查找或创建 transition 元素
            transition = sld.find(_make_p_tag('transition'), NSMAP)
            if transition is None:
                transition = etree.SubElement(sld, _make_p_tag('transition'))

            # 设置切换时间（转换为 1/1000 秒）
            transition.set('spd', 'med')  # slow, med, fast

            # 根据类型添加切换效果元素
            trans_type = trans_config.trans_type

            # 清除现有的切换效果子元素
            for child in list(transition):
                transition.remove(child)

            if trans_type == "fade":
                fade = etree.SubElement(transition, _make_p_tag('fade'))
                fade.set('thruBlk', 'false')
            elif trans_type == "push":
                push = etree.SubElement(transition, _make_p_tag('push'))
                push.set('dir', 'r')  # l, r, u, d
            elif trans_type == "wipe":
                wipe = etree.SubElement(transition, _make_p_tag('wipe'))
                wipe.set('dir', 'r')
            elif trans_type == "split":
                split = etree.SubElement(transition, _make_p_tag('split'))
                split.set('orient', 'horz')
                split.set('dir', 'out')
            elif trans_type == "reveal":
                # reveal 不是标准 OOXML，使用 fade 替代
                fade = etree.SubElement(transition, _make_p_tag('fade'))
                fade.set('thruBlk', 'true')
            elif trans_type == "random":
                random_trans = etree.SubElement(transition, _make_p_tag('random'))

            return True

        except Exception as e:
            logger.warning(f"添加页面切换效果失败: {e}")
            return False

    def apply_default_animations(
        self,
        slide: Slide,
        slide_type: str,
        shapes_map: Optional[Dict[str, List[BaseShape]]] = None
    ) -> Dict[str, int]:
        """应用默认动画配置

        Args:
            slide: 幻灯片对象
            slide_type: 页面类型（bullets, image_with_text 等）
            shapes_map: 形状映射 {"title": [shape], "items": [shape1, shape2, ...]}

        Returns:
            {"animations": 添加的动画数, "transition": 1/0}
        """
        result = {"animations": 0, "transition": 0}

        if not self.enabled:
            return result

        # 应用页面切换
        trans_config = DEFAULT_TRANSITIONS.get(slide_type)
        if trans_config and self.add_slide_transition(slide, trans_config):
            result["transition"] = 1

        # 如果没有提供形状映射，只添加切换效果
        if not shapes_map:
            return result

        # 获取该页面类型的默认动画配置
        anim_defaults = DEFAULT_SLIDE_ANIMATIONS.get(slide_type, {})

        for element_name, shapes in shapes_map.items():
            if not shapes:
                continue

            config = anim_defaults.get(element_name)
            if config is None:
                continue

            if isinstance(config, AnimationConfig):
                # 单个形状
                if isinstance(shapes, list):
                    for shape in shapes:
                        if self.add_entrance_animation(slide, shape, config):
                            result["animations"] += 1
                else:
                    if self.add_entrance_animation(slide, shapes, config):
                        result["animations"] += 1
            elif element_name == "items" or element_name.endswith("_items"):
                # 列表项使用顺序动画
                interval = anim_defaults.get("item_interval", 150)
                item_config = anim_defaults.get("items", AnimationConfig("fly_in", 300, 200))
                if isinstance(shapes, list):
                    result["animations"] += self.add_sequence_animation(
                        slide, shapes, item_config, interval
                    )

        return result

    def _get_or_create_timing(self, slide: Slide):
        """获取或创建 timing 元素"""
        sld = slide._element
        timing = sld.find(_make_p_tag('timing'), NSMAP)
        if timing is None:
            timing = etree.SubElement(sld, _make_p_tag('timing'))
        return timing

    def _get_or_create_tn_lst(self, timing):
        """获取或创建时间节点列表"""
        tn_lst = timing.find(_make_p_tag('tnLst'), NSMAP)
        if tn_lst is None:
            tn_lst = etree.SubElement(timing, _make_p_tag('tnLst'))
        return tn_lst

    def _create_animation_node(
        self,
        parent,
        shape: BaseShape,
        config: AnimationConfig
    ):
        """创建动画节点

        这是简化的实现，创建基本的动画结构。
        完整的 OOXML 动画非常复杂，这里只实现核心功能。
        """
        # 获取形状 ID
        shape_id = shape.shape_id

        # 创建 par (并行) 节点
        par = etree.SubElement(parent, _make_p_tag('par'))

        # 创建 cTn (通用时间节点)
        ctn = etree.SubElement(par, _make_p_tag('cTn'))
        ctn.set('id', str(self._shape_counter + 1))
        ctn.set('dur', str(config.duration))
        ctn.set('fill', 'hold')

        if config.delay > 0:
            ctn.set('delay', str(config.delay))

        # 创建 stCondLst (开始条件列表)
        st_cond_lst = etree.SubElement(ctn, _make_p_tag('stCondLst'))
        cond = etree.SubElement(st_cond_lst, _make_p_tag('cond'))
        cond.set('evt', 'onBegin')
        cond.set('delay', '0')

        # 创建 childTnLst (子时间节点列表)
        child_tn_lst = etree.SubElement(ctn, _make_p_tag('childTnLst'))

        # 根据动画类型创建效果
        self._create_effect_node(child_tn_lst, shape_id, config)

    def _create_effect_node(
        self,
        parent,
        shape_id: int,
        config: AnimationConfig
    ):
        """创建效果节点"""
        anim_type = config.anim_type

        # 创建 set 节点（用于控制可见性）
        set_node = etree.SubElement(parent, _make_p_tag('set'))

        cbn = etree.SubElement(set_node, _make_p_tag('cBhvr'))
        ctn = etree.SubElement(cbn, _make_p_tag('cTn'))
        ctn.set('id', str(self._shape_counter + 100))
        ctn.set('dur', '1')
        ctn.set('fill', 'hold')

        tgt_el = etree.SubElement(cbn, _make_p_tag('tgtEl'))
        sp_tgt = etree.SubElement(tgt_el, _make_p_tag('spTgt'))
        sp_tgt.set('spid', str(shape_id))

        # 根据动画类型添加特定效果
        if anim_type == "fade":
            self._add_fade_effect(parent, shape_id, config)
        elif anim_type == "fly_in":
            self._add_fly_in_effect(parent, shape_id, config)
        elif anim_type == "zoom":
            self._add_zoom_effect(parent, shape_id, config)
        elif anim_type == "wipe":
            self._add_wipe_effect(parent, shape_id, config)
        elif anim_type == "appear":
            # appear 不需要额外效果，只需要 set 节点
            pass
        elif anim_type == "bounce":
            self._add_bounce_effect(parent, shape_id, config)
        elif anim_type == "float":
            self._add_float_effect(parent, shape_id, config)

    def _add_fade_effect(self, parent, shape_id: int, config: AnimationConfig):
        """添加淡入效果"""
        anim_effect = etree.SubElement(parent, _make_p_tag('animEffect'))
        anim_effect.set('transition', 'in')
        anim_effect.set('filter', 'fade')

        cbn = etree.SubElement(anim_effect, _make_p_tag('cBhvr'))
        ctn = etree.SubElement(cbn, _make_p_tag('cTn'))
        ctn.set('id', str(self._shape_counter + 200))
        ctn.set('dur', str(config.duration))

        tgt_el = etree.SubElement(cbn, _make_p_tag('tgtEl'))
        sp_tgt = etree.SubElement(tgt_el, _make_p_tag('spTgt'))
        sp_tgt.set('spid', str(shape_id))

    def _add_fly_in_effect(self, parent, shape_id: int, config: AnimationConfig):
        """添加飞入效果"""
        # 使用 anim 节点控制位置动画
        anim = etree.SubElement(parent, _make_p_tag('anim'))
        anim.set('calcmode', 'lin')
        anim.set('valueType', 'num')

        cbn = etree.SubElement(anim, _make_p_tag('cBhvr'))
        ctn = etree.SubElement(cbn, _make_p_tag('cTn'))
        ctn.set('id', str(self._shape_counter + 200))
        ctn.set('dur', str(config.duration))
        ctn.set('fill', 'hold')

        tgt_el = etree.SubElement(cbn, _make_p_tag('tgtEl'))
        sp_tgt = etree.SubElement(tgt_el, _make_p_tag('spTgt'))
        sp_tgt.set('spid', str(shape_id))

        # 设置属性名称（位置）
        attr_name = etree.SubElement(cbn, _make_p_tag('attrName'))

        # 根据方向设置动画属性
        direction = config.direction
        if direction in ["left", "right"]:
            attr_name.text = "ppt_x"
        else:
            attr_name.text = "ppt_y"

        # 添加动画值
        tav_lst = etree.SubElement(anim, _make_p_tag('tavLst'))

        # 起始值
        tav1 = etree.SubElement(tav_lst, _make_p_tag('tav'))
        tav1.set('tm', '0')
        val1 = etree.SubElement(tav1, _make_p_tag('val'))
        str_val1 = etree.SubElement(val1, _make_p_tag('strVal'))

        # 根据方向设置起始位置
        if direction == "left":
            str_val1.set('val', '#ppt_x-1')
        elif direction == "right":
            str_val1.set('val', '#ppt_x+1')
        elif direction == "top":
            str_val1.set('val', '#ppt_y-1')
        else:  # bottom
            str_val1.set('val', '#ppt_y+1')

        # 结束值
        tav2 = etree.SubElement(tav_lst, _make_p_tag('tav'))
        tav2.set('tm', '100000')
        val2 = etree.SubElement(tav2, _make_p_tag('val'))
        str_val2 = etree.SubElement(val2, _make_p_tag('strVal'))
        str_val2.set('val', '#ppt_x' if direction in ["left", "right"] else '#ppt_y')

    def _add_zoom_effect(self, parent, shape_id: int, config: AnimationConfig):
        """添加缩放效果"""
        anim_effect = etree.SubElement(parent, _make_p_tag('animEffect'))
        anim_effect.set('transition', 'in')
        anim_effect.set('filter', 'zoom')

        cbn = etree.SubElement(anim_effect, _make_p_tag('cBhvr'))
        ctn = etree.SubElement(cbn, _make_p_tag('cTn'))
        ctn.set('id', str(self._shape_counter + 200))
        ctn.set('dur', str(config.duration))

        tgt_el = etree.SubElement(cbn, _make_p_tag('tgtEl'))
        sp_tgt = etree.SubElement(tgt_el, _make_p_tag('spTgt'))
        sp_tgt.set('spid', str(shape_id))

    def _add_wipe_effect(self, parent, shape_id: int, config: AnimationConfig):
        """添加擦除效果"""
        anim_effect = etree.SubElement(parent, _make_p_tag('animEffect'))
        anim_effect.set('transition', 'in')
        anim_effect.set('filter', 'wipe(right)')  # 可以是 left, right, up, down

        cbn = etree.SubElement(anim_effect, _make_p_tag('cBhvr'))
        ctn = etree.SubElement(cbn, _make_p_tag('cTn'))
        ctn.set('id', str(self._shape_counter + 200))
        ctn.set('dur', str(config.duration))

        tgt_el = etree.SubElement(cbn, _make_p_tag('tgtEl'))
        sp_tgt = etree.SubElement(tgt_el, _make_p_tag('spTgt'))
        sp_tgt.set('spid', str(shape_id))

    def _add_bounce_effect(self, parent, shape_id: int, config: AnimationConfig):
        """添加弹跳效果（简化实现）"""
        # 使用 fade + zoom 模拟弹跳
        self._add_zoom_effect(parent, shape_id, config)

    def _add_float_effect(self, parent, shape_id: int, config: AnimationConfig):
        """添加浮动效果（简化实现）"""
        # 使用 fade 模拟浮动
        self._add_fade_effect(parent, shape_id, config)


# 全局动画构建器实例
_animation_builder: Optional[AnimationBuilder] = None


def get_animation_builder(enabled: bool = True) -> AnimationBuilder:
    """获取全局动画构建器实例"""
    global _animation_builder
    if _animation_builder is None or _animation_builder.enabled != enabled:
        _animation_builder = AnimationBuilder(enabled)
    return _animation_builder


def apply_animations_to_slide(
    slide: Slide,
    slide_type: str,
    shapes_map: Optional[Dict[str, List[BaseShape]]] = None,
    enabled: bool = True
) -> Dict[str, int]:
    """便捷函数：为幻灯片应用动画

    Args:
        slide: 幻灯片对象
        slide_type: 页面类型
        shapes_map: 形状映射
        enabled: 是否启用动画

    Returns:
        {"animations": n, "transition": 0/1}
    """
    builder = get_animation_builder(enabled)
    return builder.apply_default_animations(slide, slide_type, shapes_map)


def add_simple_transition(
    slide: Slide,
    trans_type: TransitionType = "fade",
    duration: int = 700
) -> bool:
    """便捷函数：添加简单页面切换

    Args:
        slide: 幻灯片对象
        trans_type: 切换类型
        duration: 持续时间（毫秒）

    Returns:
        是否成功
    """
    builder = get_animation_builder(True)
    config = TransitionConfig(trans_type, duration)
    return builder.add_slide_transition(slide, config)
