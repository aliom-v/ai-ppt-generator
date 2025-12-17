"""创建高质量专业 PPT 模板

基于专业设计原则：
- 极简主义：大量留白，每页聚焦一个核心
- 专业配色：2-3 主色 + 中性色，高对比度
- 视觉层次：网格布局，一致对齐
- 精致装饰：渐变、几何图形、微妙阴影
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# 模板输出目录
TEMPLATES_DIR = Path(__file__).parent / "pptx_templates"

# 幻灯片尺寸 (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ==================== 专业配色方案 ====================

class ColorScheme:
    """专业配色方案"""

    # 1. 科技蓝 - 专业、信任、现代
    TECH_BLUE = {
        "primary": RGBColor(37, 99, 235),      # #2563EB 主蓝
        "secondary": RGBColor(59, 130, 246),   # #3B82F6 亮蓝
        "accent": RGBColor(14, 165, 233),      # #0EA5E9 天蓝
        "dark": RGBColor(15, 23, 42),          # #0F172A 深蓝黑
        "light": RGBColor(241, 245, 249),      # #F1F5F9 浅灰蓝
        "text": RGBColor(30, 41, 59),          # #1E293B 深色文字
        "text_light": RGBColor(100, 116, 139), # #64748B 浅色文字
        "white": RGBColor(255, 255, 255),
    }

    # 2. 优雅深色 - 高端、专业、沉稳
    ELEGANT_DARK = {
        "primary": RGBColor(139, 92, 246),     # #8B5CF6 紫色
        "secondary": RGBColor(168, 85, 247),   # #A855F7 亮紫
        "accent": RGBColor(236, 72, 153),      # #EC4899 粉色
        "dark": RGBColor(17, 24, 39),          # #111827 深黑
        "light": RGBColor(31, 41, 55),         # #1F2937 深灰
        "text": RGBColor(243, 244, 246),       # #F3F4F6 白色文字
        "text_light": RGBColor(156, 163, 175), # #9CA3AF 灰色文字
        "white": RGBColor(255, 255, 255),
    }

    # 3. 自然绿 - 环保、健康、清新
    NATURE_GREEN = {
        "primary": RGBColor(5, 150, 105),      # #059669 主绿
        "secondary": RGBColor(16, 185, 129),   # #10B981 亮绿
        "accent": RGBColor(20, 184, 166),      # #14B8A6 青绿
        "dark": RGBColor(6, 78, 59),           # #064E3B 深绿
        "light": RGBColor(236, 253, 245),      # #ECFDF5 浅绿白
        "text": RGBColor(20, 83, 45),          # #14532D 深绿文字
        "text_light": RGBColor(75, 85, 99),    # #4B5563 灰色文字
        "white": RGBColor(255, 255, 255),
    }

    # 4. 暖橙色 - 活力、创意、温暖
    WARM_ORANGE = {
        "primary": RGBColor(234, 88, 12),      # #EA580C 橙色
        "secondary": RGBColor(249, 115, 22),   # #F97316 亮橙
        "accent": RGBColor(245, 158, 11),      # #F59E0B 金黄
        "dark": RGBColor(67, 20, 7),           # #431407 深棕
        "light": RGBColor(255, 247, 237),      # #FFF7ED 暖白
        "text": RGBColor(55, 48, 44),          # #37302C 深色文字
        "text_light": RGBColor(120, 113, 108), # #78716C 浅色文字
        "white": RGBColor(255, 255, 255),
    }

    # 5. 极简黑白 - 简约、高端、专业
    MINIMAL_BW = {
        "primary": RGBColor(0, 0, 0),          # #000000 纯黑
        "secondary": RGBColor(64, 64, 64),     # #404040 深灰
        "accent": RGBColor(239, 68, 68),       # #EF4444 点缀红
        "dark": RGBColor(23, 23, 23),          # #171717 深黑
        "light": RGBColor(250, 250, 250),      # #FAFAFA 纯白
        "text": RGBColor(23, 23, 23),          # #171717 黑色文字
        "text_light": RGBColor(115, 115, 115), # #737373 灰色文字
        "white": RGBColor(255, 255, 255),
    }

    # 6. 商务蓝灰 - 稳重、专业、可信
    CORPORATE = {
        "primary": RGBColor(30, 64, 175),      # #1E40AF 深蓝
        "secondary": RGBColor(59, 130, 246),   # #3B82F6 亮蓝
        "accent": RGBColor(251, 191, 36),      # #FBBF24 金色
        "dark": RGBColor(30, 41, 59),          # #1E293B 深蓝灰
        "light": RGBColor(248, 250, 252),      # #F8FAFC 浅灰
        "text": RGBColor(15, 23, 42),          # #0F172A 深色文字
        "text_light": RGBColor(71, 85, 105),   # #475569 灰色文字
        "white": RGBColor(255, 255, 255),
    }


# ==================== 辅助函数 ====================

def create_presentation():
    """创建标准 16:9 演示文稿"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_shape(slide, shape_type, left, top, width, height, fill_color=None,
              line_color=None, line_width=0):
    """添加形状（统一接口）"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color and line_width > 0:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()

    return shape


def add_text_box(slide, left, top, width, height, text="", font_size=14,
                 font_color=None, bold=False, align=PP_ALIGN.LEFT,
                 vertical=MSO_ANCHOR.TOP):
    """添加文本框（统一接口）"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = False

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.bold = bold

    if font_color:
        p.font.color.rgb = font_color

    # 设置垂直对齐
    tf.paragraphs[0].vertical_anchor = vertical

    return box


def move_to_back(slide, shape):
    """将形状移到最底层（作为背景）"""
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


# ==================== 模板创建函数 ====================

def create_tech_blue_template():
    """创建科技蓝模板 - 现代专业风格"""
    prs = create_presentation()
    colors = ColorScheme.TECH_BLUE

    # ===== 封面页 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    # 深色背景
    bg = add_shape(slide1, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["dark"])
    move_to_back(slide1, bg)

    # 右上角装饰圆（大）
    add_shape(slide1, MSO_SHAPE.OVAL,
              SLIDE_WIDTH - Inches(4), Inches(-2),
              Inches(6), Inches(6), colors["primary"])

    # 右上角装饰圆（小）
    add_shape(slide1, MSO_SHAPE.OVAL,
              SLIDE_WIDTH - Inches(2.5), Inches(0.5),
              Inches(3), Inches(3), colors["secondary"])

    # 左下角装饰
    add_shape(slide1, MSO_SHAPE.OVAL,
              Inches(-1), SLIDE_HEIGHT - Inches(3),
              Inches(4), Inches(4), colors["accent"])

    # 底部渐变条
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              0, SLIDE_HEIGHT - Inches(0.08),
              SLIDE_WIDTH, Inches(0.08), colors["accent"])

    # 标题区域
    add_text_box(slide1, Inches(0.8), Inches(2.5), Inches(8), Inches(1.5),
                 "演示文稿标题", 56, colors["white"], bold=True)

    # 副标题
    add_text_box(slide1, Inches(0.8), Inches(4.2), Inches(8), Inches(0.8),
                 "专业 · 创新 · 高效", 24, colors["text_light"])

    # 日期/作者
    add_text_box(slide1, Inches(0.8), Inches(6.2), Inches(4), Inches(0.5),
                 "2024", 16, colors["text_light"])

    # ===== 内容页 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # 白色背景
    bg2 = add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["light"])
    move_to_back(slide2, bg2)

    # 顶部装饰条
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              0, 0, SLIDE_WIDTH, Inches(0.06), colors["primary"])

    # 左侧色块装饰
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              0, Inches(0.06), Inches(0.5), SLIDE_HEIGHT - Inches(0.06), colors["dark"])

    # 标题背景
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              Inches(0.5), Inches(0.06), Inches(12.833), Inches(1.2), colors["white"])

    # 标题装饰线
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              Inches(0.8), Inches(1.1), Inches(0.8), Inches(0.05), colors["primary"])

    # 标题
    add_text_box(slide2, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
                 "章节标题", 36, colors["text"], bold=True)

    # 内容区域背景
    add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.5),
              colors["white"], colors["light"], 1)

    # 右下角装饰
    add_shape(slide2, MSO_SHAPE.OVAL,
              SLIDE_WIDTH - Inches(2), SLIDE_HEIGHT - Inches(2),
              Inches(3), Inches(3), colors["accent"])

    # ===== 结束页 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # 深色背景
    bg3 = add_shape(slide3, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["dark"])
    move_to_back(slide3, bg3)

    # 中央装饰圆环
    add_shape(slide3, MSO_SHAPE.DONUT,
              Inches(5), Inches(1.5), Inches(3.333), Inches(3.333), colors["primary"])

    # 感谢文字
    add_text_box(slide3, Inches(0), Inches(5), SLIDE_WIDTH, Inches(1),
                 "感谢聆听", 48, colors["white"], bold=True, align=PP_ALIGN.CENTER)

    # 副文字
    add_text_box(slide3, Inches(0), Inches(6), SLIDE_WIDTH, Inches(0.6),
                 "THANK YOU", 20, colors["text_light"], align=PP_ALIGN.CENTER)

    # 保存
    prs.save(str(TEMPLATES_DIR / "premium_tech_blue.pptx"))
    print("✓ 科技蓝模板: premium_tech_blue.pptx")


def create_elegant_dark_template():
    """创建优雅深色模板 - 高端大气风格"""
    prs = create_presentation()
    colors = ColorScheme.ELEGANT_DARK

    # ===== 封面页 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    # 深色背景
    bg = add_shape(slide1, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["dark"])
    move_to_back(slide1, bg)

    # 渐变装饰块（左上）
    add_shape(slide1, MSO_SHAPE.PARALLELOGRAM,
              Inches(-2), Inches(-1), Inches(8), Inches(5), colors["primary"])

    # 装饰线条
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              Inches(5), Inches(3.5), Inches(7), Pt(2), colors["accent"])

    # 小装饰点
    for i in range(5):
        add_shape(slide1, MSO_SHAPE.OVAL,
                  Inches(5.5 + i * 0.4), Inches(3.4), Inches(0.15), Inches(0.15), colors["secondary"])

    # 标题
    add_text_box(slide1, Inches(5), Inches(2.2), Inches(7.5), Inches(1.2),
                 "演示标题", 60, colors["text"], bold=True)

    # 副标题
    add_text_box(slide1, Inches(5), Inches(3.8), Inches(7), Inches(0.6),
                 "优雅 · 专业 · 创意", 22, colors["text_light"])

    # 底部装饰
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              0, SLIDE_HEIGHT - Inches(0.5), SLIDE_WIDTH, Inches(0.5), colors["light"])

    # ===== 内容页 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # 深色背景
    bg2 = add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["dark"])
    move_to_back(slide2, bg2)

    # 顶部渐变条
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              0, 0, SLIDE_WIDTH, Inches(0.08), colors["primary"])
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              0, Inches(0.08), SLIDE_WIDTH * 0.7, Inches(0.04), colors["accent"])

    # 标题区域
    add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                 "内容标题", 38, colors["text"], bold=True)

    # 标题下划线
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              Inches(0.8), Inches(1.3), Inches(1.5), Pt(3), colors["primary"])

    # 内容卡片背景
    add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(0.6), Inches(1.6), Inches(12.133), Inches(5.4), colors["light"])

    # 左侧装饰
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              Inches(0.6), Inches(1.6), Inches(0.08), Inches(5.4), colors["primary"])

    # ===== 结束页 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    bg3 = add_shape(slide3, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["dark"])
    move_to_back(slide3, bg3)

    # 中央大圆
    add_shape(slide3, MSO_SHAPE.OVAL,
              Inches(4.5), Inches(1), Inches(4.333), Inches(4.333), colors["primary"])

    # 内圆
    add_shape(slide3, MSO_SHAPE.OVAL,
              Inches(5.2), Inches(1.7), Inches(2.933), Inches(2.933), colors["dark"])

    # 感谢文字
    add_text_box(slide3, Inches(0), Inches(5.5), SLIDE_WIDTH, Inches(1),
                 "THANKS", 52, colors["text"], bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide3, Inches(0), Inches(6.4), SLIDE_WIDTH, Inches(0.5),
                 "期待与您的下次交流", 18, colors["text_light"], align=PP_ALIGN.CENTER)

    prs.save(str(TEMPLATES_DIR / "premium_elegant_dark.pptx"))
    print("✓ 优雅深色模板: premium_elegant_dark.pptx")


def create_nature_green_template():
    """创建自然绿模板 - 清新环保风格"""
    prs = create_presentation()
    colors = ColorScheme.NATURE_GREEN

    # ===== 封面页 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    # 浅色背景
    bg = add_shape(slide1, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["light"])
    move_to_back(slide1, bg)

    # 右侧大色块
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              SLIDE_WIDTH - Inches(5), 0, Inches(5), SLIDE_HEIGHT, colors["primary"])

    # 装饰圆形
    add_shape(slide1, MSO_SHAPE.OVAL,
              SLIDE_WIDTH - Inches(6), Inches(1), Inches(4), Inches(4), colors["secondary"])

    add_shape(slide1, MSO_SHAPE.OVAL,
              SLIDE_WIDTH - Inches(3), Inches(4), Inches(2.5), Inches(2.5), colors["accent"])

    # 左侧内容区
    add_text_box(slide1, Inches(0.8), Inches(2.5), Inches(7), Inches(1.2),
                 "演示文稿", 58, colors["text"], bold=True)

    add_text_box(slide1, Inches(0.8), Inches(3.9), Inches(6), Inches(0.8),
                 "自然 · 清新 · 可持续", 22, colors["text_light"])

    # 底部线条
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              Inches(0.8), Inches(5), Inches(3), Pt(3), colors["primary"])

    # ===== 内容页 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    bg2 = add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["white"])
    move_to_back(slide2, bg2)

    # 顶部色条
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              0, 0, SLIDE_WIDTH, Inches(1), colors["primary"])

    # 标题（白色）
    add_text_box(slide2, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
                 "章节标题", 32, colors["white"], bold=True)

    # 左侧装饰条
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              Inches(0.6), Inches(1.3), Inches(0.06), Inches(5.7), colors["secondary"])

    # 右上角装饰
    add_shape(slide2, MSO_SHAPE.OVAL,
              SLIDE_WIDTH - Inches(2.5), Inches(-0.5), Inches(3), Inches(3), colors["accent"])

    # ===== 结束页 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    bg3 = add_shape(slide3, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["primary"])
    move_to_back(slide3, bg3)

    # 装饰圆
    add_shape(slide3, MSO_SHAPE.OVAL,
              Inches(1), Inches(1), Inches(5), Inches(5), colors["secondary"])

    add_shape(slide3, MSO_SHAPE.OVAL,
              Inches(8), Inches(3), Inches(4), Inches(4), colors["accent"])

    # 感谢文字
    add_text_box(slide3, Inches(0), Inches(3), SLIDE_WIDTH, Inches(1),
                 "感谢观看", 52, colors["white"], bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide3, Inches(0), Inches(4.2), SLIDE_WIDTH, Inches(0.5),
                 "GREEN · FRESH · NATURAL", 18, colors["light"], align=PP_ALIGN.CENTER)

    prs.save(str(TEMPLATES_DIR / "premium_nature_green.pptx"))
    print("✓ 自然绿模板: premium_nature_green.pptx")


def create_warm_orange_template():
    """创建暖橙色模板 - 活力创意风格"""
    prs = create_presentation()
    colors = ColorScheme.WARM_ORANGE

    # ===== 封面页 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    bg = add_shape(slide1, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["light"])
    move_to_back(slide1, bg)

    # 斜向色块
    add_shape(slide1, MSO_SHAPE.PARALLELOGRAM,
              Inches(-1), SLIDE_HEIGHT - Inches(4), Inches(10), Inches(5), colors["primary"])

    # 装饰圆
    add_shape(slide1, MSO_SHAPE.OVAL,
              Inches(7), Inches(0.5), Inches(5), Inches(5), colors["secondary"])

    add_shape(slide1, MSO_SHAPE.OVAL,
              Inches(10), Inches(3), Inches(3), Inches(3), colors["accent"])

    # 标题
    add_text_box(slide1, Inches(0.8), Inches(1.5), Inches(8), Inches(1.2),
                 "创意演示", 60, colors["text"], bold=True)

    add_text_box(slide1, Inches(0.8), Inches(2.9), Inches(6), Inches(0.7),
                 "活力 · 温暖 · 创新", 24, colors["text_light"])

    # ===== 内容页 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    bg2 = add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["white"])
    move_to_back(slide2, bg2)

    # 顶部装饰
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              0, 0, SLIDE_WIDTH, Inches(0.1), colors["primary"])

    # 左侧色条
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              0, 0, Inches(0.4), SLIDE_HEIGHT, colors["secondary"])

    # 标题
    add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(10), Inches(0.9),
                 "内容标题", 36, colors["text"], bold=True)

    # 标题装饰
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              Inches(0.8), Inches(1.2), Inches(1.2), Pt(4), colors["primary"])

    # 右下角装饰
    add_shape(slide2, MSO_SHAPE.OVAL,
              SLIDE_WIDTH - Inches(3), SLIDE_HEIGHT - Inches(2.5),
              Inches(4), Inches(4), colors["accent"])

    # ===== 结束页 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    bg3 = add_shape(slide3, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["primary"])
    move_to_back(slide3, bg3)

    # 装饰
    add_shape(slide3, MSO_SHAPE.OVAL,
              Inches(-1), Inches(-1), Inches(4), Inches(4), colors["secondary"])

    add_shape(slide3, MSO_SHAPE.OVAL,
              SLIDE_WIDTH - Inches(3), SLIDE_HEIGHT - Inches(3),
              Inches(4), Inches(4), colors["accent"])

    add_text_box(slide3, Inches(0), Inches(3), SLIDE_WIDTH, Inches(1),
                 "THANK YOU", 56, colors["white"], bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide3, Inches(0), Inches(4.2), SLIDE_WIDTH, Inches(0.5),
                 "感谢您的聆听", 20, colors["light"], align=PP_ALIGN.CENTER)

    prs.save(str(TEMPLATES_DIR / "premium_warm_orange.pptx"))
    print("✓ 暖橙色模板: premium_warm_orange.pptx")


def create_minimal_bw_template():
    """创建极简黑白模板 - 高端简约风格"""
    prs = create_presentation()
    colors = ColorScheme.MINIMAL_BW

    # ===== 封面页 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    bg = add_shape(slide1, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["light"])
    move_to_back(slide1, bg)

    # 大号标题
    add_text_box(slide1, Inches(0.8), Inches(2.8), Inches(11), Inches(1.5),
                 "极简主义", 72, colors["primary"], bold=True)

    # 下划线
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              Inches(0.8), Inches(4.5), Inches(2), Pt(4), colors["primary"])

    # 副标题
    add_text_box(slide1, Inches(0.8), Inches(5), Inches(8), Inches(0.6),
                 "Less is More", 24, colors["text_light"])

    # 右侧装饰线
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              SLIDE_WIDTH - Inches(0.5), Inches(1), Pt(2), Inches(5.5), colors["secondary"])

    # 点缀色块
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              SLIDE_WIDTH - Inches(0.7), Inches(1), Inches(0.4), Inches(0.4), colors["accent"])

    # ===== 内容页 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    bg2 = add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["light"])
    move_to_back(slide2, bg2)

    # 左侧粗线
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              Inches(0.6), Inches(0.5), Pt(6), Inches(1), colors["primary"])

    # 标题
    add_text_box(slide2, Inches(0.9), Inches(0.5), Inches(10), Inches(0.9),
                 "标题", 40, colors["text"], bold=True)

    # 分隔线
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              Inches(0.6), Inches(1.5), Inches(12.133), Pt(1), colors["secondary"])

    # 底部装饰
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              0, SLIDE_HEIGHT - Inches(0.3), SLIDE_WIDTH, Inches(0.3), colors["dark"])

    # ===== 结束页 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    bg3 = add_shape(slide3, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["dark"])
    move_to_back(slide3, bg3)

    # 大号 THANKS
    add_text_box(slide3, Inches(0), Inches(2.8), SLIDE_WIDTH, Inches(1.5),
                 "THANKS", 80, colors["light"], bold=True, align=PP_ALIGN.CENTER)

    # 红色点缀
    add_shape(slide3, MSO_SHAPE.RECTANGLE,
              Inches(5.9), Inches(4.5), Inches(1.5), Pt(4), colors["accent"])

    add_text_box(slide3, Inches(0), Inches(5), SLIDE_WIDTH, Inches(0.5),
                 "简约不简单", 18, colors["text_light"], align=PP_ALIGN.CENTER)

    prs.save(str(TEMPLATES_DIR / "premium_minimal_bw.pptx"))
    print("✓ 极简黑白模板: premium_minimal_bw.pptx")


def create_corporate_template():
    """创建商务蓝灰模板 - 专业稳重风格"""
    prs = create_presentation()
    colors = ColorScheme.CORPORATE

    # ===== 封面页 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    # 上半部深色
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              0, 0, SLIDE_WIDTH, SLIDE_HEIGHT * 0.65, colors["dark"])

    # 下半部浅色
    bg = add_shape(slide1, MSO_SHAPE.RECTANGLE,
              0, SLIDE_HEIGHT * 0.65, SLIDE_WIDTH, SLIDE_HEIGHT * 0.35, colors["light"])
    move_to_back(slide1, bg)

    # 金色装饰线
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              0, SLIDE_HEIGHT * 0.65 - Inches(0.05), SLIDE_WIDTH, Inches(0.1), colors["accent"])

    # 左侧蓝色装饰
    add_shape(slide1, MSO_SHAPE.RECTANGLE,
              0, 0, Inches(0.15), SLIDE_HEIGHT * 0.65, colors["primary"])

    # 标题
    add_text_box(slide1, Inches(0.8), Inches(1.8), Inches(10), Inches(1.2),
                 "企业报告", 54, colors["white"], bold=True)

    add_text_box(slide1, Inches(0.8), Inches(3.2), Inches(8), Inches(0.7),
                 "CORPORATE PRESENTATION", 20, colors["text_light"])

    # 底部信息
    add_text_box(slide1, Inches(0.8), Inches(5.5), Inches(4), Inches(0.5),
                 "2024 年度报告", 16, colors["text"])

    # ===== 内容页 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    bg2 = add_shape(slide2, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["light"])
    move_to_back(slide2, bg2)

    # 顶部深色条
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              0, 0, SLIDE_WIDTH, Inches(1.2), colors["dark"])

    # 金色分隔线
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              0, Inches(1.2), SLIDE_WIDTH, Inches(0.06), colors["accent"])

    # 左侧蓝色标记
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              0, 0, Inches(0.1), Inches(1.2), colors["primary"])

    # 标题
    add_text_box(slide2, Inches(0.6), Inches(0.35), Inches(10), Inches(0.7),
                 "章节标题", 32, colors["white"], bold=True)

    # 右上角装饰
    add_shape(slide2, MSO_SHAPE.RECTANGLE,
              SLIDE_WIDTH - Inches(1.5), 0, Inches(1.5), Inches(1.2), colors["primary"])

    # ===== 结束页 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    bg3 = add_shape(slide3, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, colors["dark"])
    move_to_back(slide3, bg3)

    # 金色大圆
    add_shape(slide3, MSO_SHAPE.OVAL,
              Inches(4.5), Inches(1.5), Inches(4.333), Inches(4.333), colors["accent"])

    # 内圆
    add_shape(slide3, MSO_SHAPE.OVAL,
              Inches(5.3), Inches(2.3), Inches(2.733), Inches(2.733), colors["dark"])

    # 感谢
    add_text_box(slide3, Inches(0), Inches(6), SLIDE_WIDTH, Inches(0.8),
                 "THANK YOU FOR YOUR ATTENTION", 24, colors["white"], bold=True, align=PP_ALIGN.CENTER)

    # 底部蓝色线
    add_shape(slide3, MSO_SHAPE.RECTANGLE,
              Inches(4), Inches(6.8), Inches(5.333), Pt(3), colors["primary"])

    prs.save(str(TEMPLATES_DIR / "premium_corporate.pptx"))
    print("✓ 商务蓝灰模板: premium_corporate.pptx")


def create_all_premium_templates():
    """创建所有高质量模板"""
    print("\n" + "=" * 60)
    print("创建高质量专业 PPT 模板")
    print("=" * 60)

    TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)

    create_tech_blue_template()
    create_elegant_dark_template()
    create_nature_green_template()
    create_warm_orange_template()
    create_minimal_bw_template()
    create_corporate_template()

    print("=" * 60)
    print(f"完成！共创建 6 个高质量模板")
    print(f"保存位置: {TEMPLATES_DIR}")
    print("=" * 60 + "\n")


if __name__ == "__main__":
    create_all_premium_templates()
