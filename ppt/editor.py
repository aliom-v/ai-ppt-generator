"""PPT 编辑器模块 - 支持已生成 PPT 的后期编辑"""
import os
import copy
from typing import Optional, List, Dict, Any
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from utils.logger import get_logger

logger = get_logger("ppt_editor")


@dataclass
class SlideInfo:
    """幻灯片信息"""
    index: int              # 幻灯片索引（0-based）
    title: str = ""         # 标题
    content_type: str = ""  # 内容类型猜测
    shape_count: int = 0    # 形状数量
    has_image: bool = False # 是否包含图片
    text_preview: str = ""  # 文本预览


class PPTEditor:
    """PPT 编辑器

    支持对已生成的 PPT 进行编辑操作：
    - 修改页面标题和内容
    - 调整页面顺序
    - 删除页面
    - 复制页面
    """

    def __init__(self, ppt_path: str):
        """初始化编辑器

        Args:
            ppt_path: PPT 文件路径
        """
        if not os.path.exists(ppt_path):
            raise FileNotFoundError(f"PPT 文件不存在: {ppt_path}")

        self.ppt_path = ppt_path
        self.prs = Presentation(ppt_path)
        self._modified = False

    def get_slide_count(self) -> int:
        """获取幻灯片数量"""
        return len(self.prs.slides)

    def get_slide_info(self, index: int) -> Optional[SlideInfo]:
        """获取指定幻灯片的信息

        Args:
            index: 幻灯片索引（0-based）

        Returns:
            SlideInfo 或 None
        """
        if index < 0 or index >= len(self.prs.slides):
            return None

        slide = self.prs.slides[index]
        info = SlideInfo(index=index)

        # 提取标题
        if slide.shapes.title:
            info.title = slide.shapes.title.text or ""

        # 统计形状和检测内容类型
        text_parts = []
        for shape in slide.shapes:
            info.shape_count += 1

            # 检测图片
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                info.has_image = True

            # 提取文本
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text[:100])

        # 生成文本预览
        if text_parts:
            info.text_preview = " | ".join(text_parts)[:200]

        # 猜测内容类型
        info.content_type = self._guess_content_type(slide, info)

        return info

    def _guess_content_type(self, slide, info: SlideInfo) -> str:
        """猜测幻灯片内容类型"""
        if info.has_image:
            return "image_with_text"

        # 根据形状数量和布局猜测
        if info.shape_count <= 3:
            if "谢谢" in info.title or "Thank" in info.title.lower():
                return "ending"
            return "quote"

        if info.shape_count > 8:
            return "timeline"

        return "bullets"

    def list_slides(self) -> List[SlideInfo]:
        """列出所有幻灯片信息"""
        return [self.get_slide_info(i) for i in range(len(self.prs.slides))]

    def update_slide_title(self, index: int, new_title: str) -> bool:
        """更新幻灯片标题

        Args:
            index: 幻灯片索引
            new_title: 新标题

        Returns:
            是否成功
        """
        if index < 0 or index >= len(self.prs.slides):
            return False

        slide = self.prs.slides[index]
        if slide.shapes.title:
            slide.shapes.title.text = new_title
            self._modified = True
            logger.info(f"更新幻灯片 {index} 标题: {new_title}")
            return True

        return False

    def update_text_content(self, index: int, shape_index: int, new_text: str) -> bool:
        """更新幻灯片中指定形状的文本内容

        Args:
            index: 幻灯片索引
            shape_index: 形状索引
            new_text: 新文本

        Returns:
            是否成功
        """
        if index < 0 or index >= len(self.prs.slides):
            return False

        slide = self.prs.slides[index]
        shapes_with_text = [s for s in slide.shapes if hasattr(s, "text_frame")]

        if shape_index < 0 or shape_index >= len(shapes_with_text):
            return False

        shape = shapes_with_text[shape_index]
        # 保留格式，只更新文本
        if shape.text_frame.paragraphs:
            shape.text_frame.paragraphs[0].text = new_text
            self._modified = True
            logger.info(f"更新幻灯片 {index} 形状 {shape_index} 文本")
            return True

        return False

    def delete_slide(self, index: int) -> bool:
        """删除幻灯片

        Args:
            index: 要删除的幻灯片索引

        Returns:
            是否成功
        """
        if index < 0 or index >= len(self.prs.slides):
            return False

        # python-pptx 删除幻灯片需要操作底层 XML
        slide_id = self.prs.slides._sldIdLst[index]
        rId = slide_id.rId
        self.prs.part.drop_rel(rId)
        del self.prs.slides._sldIdLst[index]

        self._modified = True
        logger.info(f"删除幻灯片 {index}")
        return True

    def move_slide(self, from_index: int, to_index: int) -> bool:
        """移动幻灯片位置

        Args:
            from_index: 原位置
            to_index: 目标位置

        Returns:
            是否成功
        """
        slide_count = len(self.prs.slides)
        if from_index < 0 or from_index >= slide_count:
            return False
        if to_index < 0 or to_index >= slide_count:
            return False
        if from_index == to_index:
            return True

        # 获取要移动的幻灯片 ID
        sldIdLst = self.prs.slides._sldIdLst
        slide_id = sldIdLst[from_index]

        # 从原位置移除
        del sldIdLst[from_index]

        # 插入到新位置
        sldIdLst.insert(to_index, slide_id)

        self._modified = True
        logger.info(f"移动幻灯片 {from_index} -> {to_index}")
        return True

    def duplicate_slide(self, index: int) -> int:
        """复制幻灯片

        Args:
            index: 要复制的幻灯片索引

        Returns:
            新幻灯片的索引，失败返回 -1
        """
        if index < 0 or index >= len(self.prs.slides):
            return -1

        try:
            # 获取原幻灯片
            source_slide = self.prs.slides[index]

            # 使用相同的布局创建新幻灯片
            slide_layout = source_slide.slide_layout
            new_slide = self.prs.slides.add_slide(slide_layout)

            # 复制所有形状（简化实现，只复制文本）
            for shape in source_slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # 尝试在新幻灯片中找到对应的形状并复制文本
                    for new_shape in new_slide.shapes:
                        if hasattr(new_shape, "text_frame"):
                            if new_shape.left == shape.left and new_shape.top == shape.top:
                                new_shape.text_frame.text = shape.text
                                break

            new_index = len(self.prs.slides) - 1
            self._modified = True
            logger.info(f"复制幻灯片 {index} -> {new_index}")
            return new_index

        except Exception as e:
            logger.error(f"复制幻灯片失败: {e}")
            return -1

    def save(self, output_path: Optional[str] = None) -> str:
        """保存修改

        Args:
            output_path: 输出路径（None 则覆盖原文件）

        Returns:
            保存的文件路径
        """
        if output_path is None:
            output_path = self.ppt_path

        # 确保目录存在
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)

        self.prs.save(output_path)
        self._modified = False
        logger.info(f"PPT 已保存: {output_path}")
        return output_path

    def save_as(self, output_path: str) -> str:
        """另存为新文件

        Args:
            output_path: 输出路径

        Returns:
            保存的文件路径
        """
        return self.save(output_path)

    @property
    def is_modified(self) -> bool:
        """是否有未保存的修改"""
        return self._modified

    def to_dict(self) -> Dict[str, Any]:
        """转换为字典（用于 API 响应）"""
        slides = []
        for i in range(len(self.prs.slides)):
            info = self.get_slide_info(i)
            if info:
                slides.append({
                    "index": info.index,
                    "title": info.title,
                    "content_type": info.content_type,
                    "shape_count": info.shape_count,
                    "has_image": info.has_image,
                    "text_preview": info.text_preview
                })

        return {
            "path": self.ppt_path,
            "slide_count": len(self.prs.slides),
            "is_modified": self._modified,
            "slides": slides
        }


def get_ppt_info(ppt_path: str) -> Optional[Dict[str, Any]]:
    """获取 PPT 文件信息（便捷函数）

    Args:
        ppt_path: PPT 文件路径

    Returns:
        PPT 信息字典
    """
    try:
        editor = PPTEditor(ppt_path)
        return editor.to_dict()
    except Exception as e:
        logger.error(f"获取 PPT 信息失败: {e}")
        return None


def edit_slide_title(ppt_path: str, index: int, new_title: str, save: bool = True) -> bool:
    """编辑幻灯片标题（便捷函数）

    Args:
        ppt_path: PPT 文件路径
        index: 幻灯片索引
        new_title: 新标题
        save: 是否保存

    Returns:
        是否成功
    """
    try:
        editor = PPTEditor(ppt_path)
        if editor.update_slide_title(index, new_title):
            if save:
                editor.save()
            return True
        return False
    except Exception as e:
        logger.error(f"编辑标题失败: {e}")
        return False


def reorder_slides(ppt_path: str, new_order: List[int], save: bool = True) -> bool:
    """重新排列幻灯片顺序（便捷函数）

    Args:
        ppt_path: PPT 文件路径
        new_order: 新顺序（索引列表）
        save: 是否保存

    Returns:
        是否成功
    """
    try:
        editor = PPTEditor(ppt_path)
        slide_count = editor.get_slide_count()

        # 验证新顺序
        if len(new_order) != slide_count:
            logger.error("新顺序长度与幻灯片数量不匹配")
            return False

        if set(new_order) != set(range(slide_count)):
            logger.error("新顺序包含无效索引")
            return False

        # 逐个移动到正确位置
        for target_pos, source_idx in enumerate(new_order):
            # 找到当前该幻灯片的位置
            current_pos = None
            for i in range(slide_count):
                info = editor.get_slide_info(i)
                if info and info.index == source_idx:
                    current_pos = i
                    break

            if current_pos is not None and current_pos != target_pos:
                editor.move_slide(current_pos, target_pos)

        if save:
            editor.save()
        return True

    except Exception as e:
        logger.error(f"重排幻灯片失败: {e}")
        return False
