"""模板预览生成模块

生成模板的预览图片，用于前端展示。
"""
import os
import hashlib
from pathlib import Path
from typing import Dict, List, Optional, Any
from dataclasses import dataclass

from pptx import Presentation
from pptx.util import Inches, Pt

from utils.logger import get_logger
from ppt.template_manager import template_manager

logger = get_logger("template_preview")


@dataclass
class TemplateInfo:
    """模板信息"""
    id: str
    name: str
    path: str
    preview_url: Optional[str] = None
    slide_count: int = 0
    width: float = 0
    height: float = 0
    layouts: List[str] = None
    colors: List[str] = None
    fonts: List[str] = None

    def to_dict(self) -> Dict:
        return {
            "id": self.id,
            "name": self.name,
            "path": self.path,
            "preview_url": self.preview_url,
            "slide_count": self.slide_count,
            "width": self.width,
            "height": self.height,
            "layouts": self.layouts or [],
            "colors": self.colors or [],
            "fonts": self.fonts or [],
        }


class TemplateAnalyzer:
    """模板分析器

    分析 PPTX 模板的结构和样式信息。
    """

    def analyze(self, template_path: str) -> TemplateInfo:
        """分析模板"""
        try:
            prs = Presentation(template_path)

            # 基本信息
            width = prs.slide_width.inches
            height = prs.slide_height.inches

            # 幻灯片母版和布局
            layouts = []
            colors = set()
            fonts = set()

            for master in prs.slide_masters:
                for layout in master.slide_layouts:
                    layouts.append(layout.name)

                # 提取主题颜色
                if hasattr(master, 'theme') and master.theme:
                    theme = master.theme
                    if hasattr(theme, 'theme_elements'):
                        # 这里可以进一步提取颜色信息
                        pass

            # 分析幻灯片
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.name:
                                    fonts.add(run.font.name)

            template_id = Path(template_path).stem
            return TemplateInfo(
                id=template_id,
                name=template_id.replace("_", " ").title(),
                path=template_path,
                slide_count=len(prs.slides),
                width=round(width, 2),
                height=round(height, 2),
                layouts=layouts[:10],  # 限制数量
                fonts=list(fonts)[:5],
            )

        except Exception as e:
            logger.error(f"分析模板失败: {template_path} - {e}")
            return TemplateInfo(
                id=Path(template_path).stem,
                name=Path(template_path).stem,
                path=template_path,
            )


class TemplatePreviewGenerator:
    """模板预览生成器

    为模板生成预览图片。
    """

    def __init__(
        self,
        preview_dir: str = "web/static/previews",
        width: int = 400,
        height: int = 300,
    ):
        self.preview_dir = preview_dir
        self.width = width
        self.height = height
        self._analyzer = TemplateAnalyzer()

        # 确保目录存在
        Path(preview_dir).mkdir(parents=True, exist_ok=True)

    def _get_preview_path(self, template_path: str) -> str:
        """获取预览图片路径"""
        # 使用模板路径的哈希作为文件名
        hash_name = hashlib.md5(template_path.encode()).hexdigest()[:8]
        template_name = Path(template_path).stem
        return os.path.join(self.preview_dir, f"{template_name}_{hash_name}.png")

    def generate_preview(self, template_path: str, force: bool = False) -> Optional[str]:
        """生成模板预览图片

        Args:
            template_path: 模板路径
            force: 是否强制重新生成

        Returns:
            预览图片路径，如果失败返回 None
        """
        preview_path = self._get_preview_path(template_path)

        # 检查是否需要重新生成
        if not force and os.path.exists(preview_path):
            template_mtime = os.path.getmtime(template_path)
            preview_mtime = os.path.getmtime(preview_path)
            if preview_mtime >= template_mtime:
                return preview_path

        try:
            from utils.export import ThumbnailExporter

            exporter = ThumbnailExporter(width=self.width, height=self.height)
            if exporter.is_available():
                return exporter.export(template_path, preview_path)

        except Exception as e:
            logger.warning(f"生成预览图片失败: {template_path} - {e}")

        # 如果导出器不可用，创建占位图
        return self._create_placeholder(template_path, preview_path)

    def _create_placeholder(self, template_path: str, preview_path: str) -> str:
        """创建占位预览图"""
        try:
            from PIL import Image, ImageDraw, ImageFont

            # 创建灰色背景
            img = Image.new('RGB', (self.width, self.height), color='#f0f0f0')
            draw = ImageDraw.Draw(img)

            # 添加模板名称
            template_name = Path(template_path).stem
            text = template_name[:20]

            # 计算文本位置（居中）
            try:
                font = ImageFont.truetype("arial.ttf", 20)
            except Exception:
                font = ImageFont.load_default()

            bbox = draw.textbbox((0, 0), text, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
            x = (self.width - text_width) // 2
            y = (self.height - text_height) // 2

            draw.text((x, y), text, fill='#666666', font=font)

            # 添加边框
            draw.rectangle(
                [(0, 0), (self.width - 1, self.height - 1)],
                outline='#cccccc',
                width=2,
            )

            img.save(preview_path)
            return preview_path

        except ImportError:
            logger.warning("PIL 未安装，无法创建占位图")
            return None

    def get_template_info(self, template_path: str) -> TemplateInfo:
        """获取模板完整信息（包括预览）"""
        info = self._analyzer.analyze(template_path)

        # 生成预览
        preview_path = self.generate_preview(template_path)
        if preview_path:
            # 转换为 URL 路径
            info.preview_url = f"/static/previews/{os.path.basename(preview_path)}"

        return info


class TemplatePreviewManager:
    """模板预览管理器

    管理所有模板的预览生成。

    用法:
        manager = TemplatePreviewManager()

        # 获取所有模板（带预览）
        templates = manager.get_all_templates()

        # 刷新预览
        manager.refresh_all_previews()
    """

    def __init__(self, preview_dir: str = "web/static/previews"):
        self._generator = TemplatePreviewGenerator(preview_dir)
        self._cache: Dict[str, TemplateInfo] = {}

    def get_template_info(self, template_id: str) -> Optional[TemplateInfo]:
        """获取单个模板信息"""
        template_path = template_manager.get_template(template_id)
        if not template_path:
            return None

        if template_id not in self._cache:
            self._cache[template_id] = self._generator.get_template_info(template_path)

        return self._cache[template_id]

    def get_all_templates(self, refresh: bool = False) -> List[Dict]:
        """获取所有模板信息"""
        templates = template_manager.list_templates()
        result = []

        for tpl in templates:
            if refresh or tpl['id'] not in self._cache:
                template_path = tpl.get('path', template_manager.get_template(tpl['id']))
                if template_path:
                    self._cache[tpl['id']] = self._generator.get_template_info(template_path)

            if tpl['id'] in self._cache:
                result.append(self._cache[tpl['id']].to_dict())
            else:
                result.append(tpl)

        return result

    def refresh_preview(self, template_id: str) -> bool:
        """刷新单个模板预览"""
        template_path = template_manager.get_template(template_id)
        if not template_path:
            return False

        self._generator.generate_preview(template_path, force=True)
        self._cache.pop(template_id, None)
        return True

    def refresh_all_previews(self):
        """刷新所有模板预览"""
        templates = template_manager.list_templates()
        for tpl in templates:
            self.refresh_preview(tpl['id'])

        logger.info(f"已刷新 {len(templates)} 个模板预览")


# 全局管理器
_preview_manager: Optional[TemplatePreviewManager] = None


def get_preview_manager() -> TemplatePreviewManager:
    """获取全局预览管理器"""
    global _preview_manager
    if _preview_manager is None:
        _preview_manager = TemplatePreviewManager()
    return _preview_manager
